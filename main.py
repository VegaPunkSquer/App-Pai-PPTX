import sys
import os

# --- RESOLUÇÃO DE CAMINHOS ABSOLUTOS E PASTA DE SISTEMA ---
if getattr(sys, 'frozen', False):
    base_dir = os.path.dirname(sys.executable)
    bundle_dir = sys._MEIPASS
else:
    base_dir = os.path.dirname(os.path.abspath(__file__))
    bundle_dir = base_dir

# O GOLPE DO APPDATA: Cria uma pasta invisível no sistema para configs e históricos.
# Assim o .exe pode ficar em qualquer lugar sem dar erro de permissão no Windows.
if os.name == 'nt':
    user_data_dir = os.path.join(os.getenv('APPDATA'), "VegaTech", "SmartSlidesPro")
else:
    user_data_dir = os.path.join(os.path.expanduser("~"), ".vegatech", "smartslidespro")

os.makedirs(user_data_dir, exist_ok=True) # Cria a pasta automaticamente e silenciosamente

# Todas as gravações do app agora vão para a pasta segura
CONFIG_FILE = os.path.join(user_data_dir, "config.json")
HISTORICO_JSON_FILE = os.path.join(user_data_dir, "historico_roteiros.json")
PASTA_HISTORICO_PPTX = os.path.join(user_data_dir, "historico_geracoes")
LOG_FILE = os.path.join(user_data_dir, "erro_ia_log.txt")
TEMP_PPTX_FILE = os.path.join(user_data_dir, "apresentacao_ia_temp.pptx")

# ==========================================
# 🚚 MIGRAÇÃO AUTOMÁTICA DE ARQUIVOS ANTIGOS
# ==========================================
import shutil

# 1. Migra as Configurações
old_config = os.path.join(base_dir, "config.json")
if os.path.exists(old_config) and not os.path.exists(CONFIG_FILE):
    try: shutil.copy2(old_config, CONFIG_FILE)
    except: pass

# 2. Migra a Lista de Histórico
old_hist_json = os.path.join(base_dir, "historico_roteiros.json")
if os.path.exists(old_hist_json) and not os.path.exists(HISTORICO_JSON_FILE):
    try: shutil.copy2(old_hist_json, HISTORICO_JSON_FILE)
    except: pass
    
# 3. Migra as Apresentações Físicas Salvas
old_hist_folder = os.path.join(base_dir, "historico_geracoes")
if os.path.exists(old_hist_folder) and not os.path.exists(PASTA_HISTORICO_PPTX):
    try: shutil.copytree(old_hist_folder, PASTA_HISTORICO_PPTX)
    except: pass

# ==========================================
# 📋 LEITURA AUTOMÁTICA DO MANIFESTO (VERSÃO E ID)
# ==========================================
MANIFESTO_FILE = os.path.join(base_dir, "vega_manifesto.json")
VERSAO_ATUAL = "v1.0.0"
PRODUTO_ID_MASTER = 1 # Usado para bater na sua API

if os.path.exists(MANIFESTO_FILE):
    try:
        import json as _json_temp
        with open(MANIFESTO_FILE, "r", encoding="utf-8") as _f:
            _dados_manifesto = _json_temp.load(_f)
            VERSAO_ATUAL = _dados_manifesto.get("versao_atual", "v1.0.0")
            if "produto_id_master" in _dados_manifesto:
                PRODUTO_ID_MASTER = _dados_manifesto["produto_id_master"]
    except Exception:
        pass

# ==========================================
# 🚀 O SEGREDO DA VELOCIDADE: SPLASH SCREEN NO TOPO
# ==========================================
# Importamos APENAS o essencial visual para a tela abrir instantaneamente
from PySide6.QtWidgets import QApplication, QSplashScreen, QLabel, QProgressBar
from PySide6.QtGui import QPixmap
from PySide6.QtCore import Qt

class SplashScreenVega(QSplashScreen):
    def __init__(self, caminho_img):
        pixmap = QPixmap(caminho_img).scaled(650, 400, Qt.KeepAspectRatio, Qt.SmoothTransformation)
        super().__init__(pixmap, Qt.FramelessWindowHint) # Removido o WindowStaysOnTopHint que travava a tela
        
        self.lbl_status = QLabel("Iniciando SmartSlides Pro...", self)
        self.lbl_status.setStyleSheet("color: white; font-weight: bold; font-size: 13px; background-color: rgba(0,0,0,150); padding: 2px; border-radius: 4px;")
        self.lbl_status.setAlignment(Qt.AlignCenter)
        self.lbl_status.setGeometry(20, pixmap.height() - 70, pixmap.width() - 40, 25)
        
        self.barra = QProgressBar(self)
        self.barra.setGeometry(20, pixmap.height() - 40, pixmap.width() - 40, 20)
        self.barra.setStyleSheet("""
            QProgressBar { border: 1px solid #555; border-radius: 5px; background-color: #222; text-align: center; color: white; font-weight: bold; }
            QProgressBar::chunk { background-color: #0078d7; border-radius: 4px; }
        """)

    def atualizar(self, valor, texto):
        self.barra.setValue(valor)
        self.lbl_status.setText(texto)
        QApplication.processEvents()

# Inicia o núcleo e mostra a imagem ANTES do Python ler o resto do arquivo!
app = None
splash = None
if __name__ == "__main__":
    app = QApplication(sys.argv)
    caminho_splash = os.path.join(bundle_dir, "assets", "SmartSlides.png")
    splash = SplashScreenVega(caminho_splash)
    splash.show()
    splash.atualizar(5, "Iniciando o núcleo do aplicativo...")

# ==========================================
# IMPORTS PESADOS (Carregam enquanto a Splash já está na tela)
# ==========================================
if splash: splash.atualizar(15, "Carregando bibliotecas de interface (PySide6)...")
import time
from PySide6.QtWidgets import (
    QDialog, QFrame, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
    QPushButton, QFileDialog, QColorDialog, QFontDialog, QMessageBox,
    QCheckBox, QSpinBox, QTextEdit, QScrollArea, QListWidget, QSlider, QProgressDialog,
    QTabWidget, QComboBox, QLineEdit, QInputDialog
)
from PySide6.QtGui import QColor, QFont, QPalette, QIcon, QPainter, QPen
from PySide6.QtCore import QThread, QSize, Signal, QPoint, QRect

if splash: splash.atualizar(40, "Carregando motor do PowerPoint (isso leva uns segundos)...")
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR

if splash: splash.atualizar(75, "Conectando motores de Inteligência Artificial...")
import requests
import motor_ia
import json
import subprocess
import traceback
import webbrowser
import ctypes
import comtypes.client
from atualizador import checar_e_atualizar

# --- JANELA DE ESCOLHA DE IMAGENS ---
class ImageSelectionDialog(QDialog):
    def __init__(self, slide_images_dict, roteiro, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Escolha as Imagens dos Slides")
        self.resize(900, 700)
        self.slide_images_dict = slide_images_dict
        self.roteiro = roteiro
        self.selected_images = {} 
        self.init_ui()

    def init_ui(self):
        layout = QHBoxLayout()
        self.list_slides = QListWidget()
        self.list_slides.setFixedWidth(200)
        self.list_slides.setStyleSheet("font-size: 14px; padding: 5px;")
        for slide_idx in self.slide_images_dict.keys():
            self.list_slides.addItem(f"Slide {slide_idx + 1}")
        
        self.list_slides.currentRowChanged.connect(self.show_options_for_slide)

        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        
        self.container_opcoes = QWidget()
        self.layout_opcoes = QVBoxLayout(self.container_opcoes)
        self.layout_opcoes.setAlignment(Qt.AlignTop)
        self.scroll_area.setWidget(self.container_opcoes)

        layout.addWidget(self.list_slides)
        layout.addWidget(self.scroll_area)

        main_layout = QVBoxLayout()
        main_layout.addLayout(layout)
        
        btn_concluir = QPushButton("Concluir Escolhas e Gerar PPTX")
        btn_concluir.setStyleSheet("background-color: #28a745; color: white; font-weight: bold; padding: 15px; font-size: 16px; border-radius: 5px;")
        btn_concluir.setCursor(Qt.PointingHandCursor)
        btn_concluir.clicked.connect(self.accept)
        main_layout.addWidget(btn_concluir)

        self.setLayout(main_layout)

        if self.list_slides.count() > 0:
            self.list_slides.setCurrentRow(0)

    def show_options_for_slide(self, row_idx):
        if row_idx < 0: return

        while self.layout_opcoes.count():
            item = self.layout_opcoes.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

        try:
            slide_idx = list(self.slide_images_dict.keys())[row_idx]
            caminhos = self.slide_images_dict[slide_idx]
        except IndexError: return

        dados_slide = self.roteiro[slide_idx]
        titulo_slide = dados_slide.get("titulo", "Sem Título")
        texto_resumo = dados_slide.get("texto", "")[:120] + "..."
        
        lbl_instrucao = QLabel(f"Slide {slide_idx + 1}: {titulo_slide}")
        lbl_instrucao.setFont(QFont("Arial", 16, weight=QFont.Bold))
        lbl_instrucao.setWordWrap(True)
        self.layout_opcoes.addWidget(lbl_instrucao)

        lbl_contexto = QLabel(f"Contexto: {texto_resumo}")
        lbl_contexto.setFont(QFont("Arial", 11, italic=True))
        lbl_contexto.setStyleSheet("padding-bottom: 15px; color: gray;")
        lbl_contexto.setWordWrap(True)
        self.layout_opcoes.addWidget(lbl_contexto)
        
        # --- LÓGICA DO CHECKBOX DE PULAR IMAGEM (ITEM 4) ---
        chk_sem_imagem = QCheckBox("❌ NÃO USAR NENHUMA IMAGEM NESTE SLIDE")
        chk_sem_imagem.setStyleSheet("font-weight: bold; color: #a8201a; font-size: 14px; padding: 5px;")
        chk_sem_imagem.setCursor(Qt.PointingHandCursor)
        
        # Se estiver marcado como "NONE" no dicionário, liga o checkbox
        if self.selected_images.get(slide_idx) == "NONE":
            chk_sem_imagem.setChecked(True)
            
        chk_sem_imagem.toggled.connect(lambda checked, s=slide_idx, r=row_idx: self.toggle_no_image(checked, s, r))
        self.layout_opcoes.addWidget(chk_sem_imagem)

        for caminho in caminhos:
            container_img = QWidget()
            vbox = QVBoxLayout(container_img)
            vbox.setAlignment(Qt.AlignCenter)

            lbl_img = QLabel()
            pixmap = QPixmap(caminho).scaled(500, 350, Qt.KeepAspectRatio, Qt.SmoothTransformation)
            lbl_img.setPixmap(pixmap)
            lbl_img.setAlignment(Qt.AlignCenter)

            btn_select = QPushButton()
            btn_select.setCursor(Qt.PointingHandCursor)
            
            if self.selected_images.get(slide_idx) == caminho:
                btn_select.setText("✅ IMAGEM SELECIONADA")
                btn_select.setStyleSheet("background-color: #28a745; color: white; font-weight: bold; padding: 12px; font-size: 14px; border-radius: 5px;")
                container_img.setStyleSheet("background-color: rgba(40, 167, 69, 0.2); border: 2px solid #28a745; border-radius: 10px;")
            else:
                btn_select.setText("Escolher esta imagem")
                btn_select.setStyleSheet("background-color: #e0e0e0; color: black; font-weight: bold; padding: 12px; font-size: 14px; border-radius: 5px;")
                container_img.setStyleSheet("border: 1px solid transparent;")

            btn_select.clicked.connect(lambda checked=False, s=slide_idx, c=caminho, r=row_idx: self.select_image(s, c, r))

            vbox.addWidget(lbl_img)
            vbox.addWidget(btn_select)
            self.layout_opcoes.addWidget(container_img)
            
        self.layout_opcoes.addStretch()

    def select_image(self, slide_idx, caminho, row_idx):
        self.selected_images[slide_idx] = caminho
        self.show_options_for_slide(row_idx)
        
    def toggle_no_image(self, checked, slide_idx, row_idx):
        if checked:
            self.selected_images[slide_idx] = "NONE" # Registra que não quer imagem
        else:
            if self.selected_images.get(slide_idx) == "NONE":
                del self.selected_images[slide_idx] # Remove a trava se desmarcar
        self.show_options_for_slide(row_idx)
        
class WorkerAuthApp(QThread):
    sucesso = Signal(str)
    erro = Signal(str)

    def __init__(self, modo, payload):
        super().__init__()
        self.modo = modo # 'login' ou 'registrar'
        self.payload = payload

    def run(self):
        url = "https://vegap-masterapp.hf.space"
        endpoint = f"{url}/auth/app/{self.modo}"
        try:
            resp = requests.post(endpoint, json=self.payload, timeout=10)
            if resp.status_code in [200, 201]:
                dados = resp.json()
                self.sucesso.emit(dados.get("access_token", ""))
            else:
                self.erro.emit(f"Erro: {resp.json().get('detail', resp.text)}")
        except Exception as e:
            self.erro.emit(f"Falha de conexão com o Master: {str(e)}")

class LoginCadastroDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Conta VegaTech")
        # Removemos o setFixedSize. Agora ela tem um tamanho mínimo confortável e pode esticar!
        self.setMinimumSize(400, 550)
        self.resize(400, 600)
        self.token_recebido = None
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout(self)
        self.tabs = QTabWidget()
        
        # ABA LOGIN
        self.tab_login = QWidget()
        layout_login = QVBoxLayout(self.tab_login)
        self.inp_login_email = QLineEdit()
        self.inp_login_email.setPlaceholderText("E-mail")
        self.inp_login_senha = QLineEdit()
        self.inp_login_senha.setPlaceholderText("Senha")
        self.inp_login_senha.setEchoMode(QLineEdit.Password)
        self.btn_logar = QPushButton("Entrar")
        self.btn_logar.setStyleSheet("background-color: #0078d7; color: white; padding: 10px; font-weight: bold;")
        self.btn_logar.clicked.connect(self.fazer_login)
        
        layout_login.addWidget(QLabel("Já tem uma conta? Acesse:"))
        layout_login.addWidget(QLabel("E-mail:"))
        layout_login.addWidget(self.inp_login_email)
        layout_login.addWidget(QLabel("Senha:"))
        layout_login.addWidget(self.inp_login_senha)
        layout_login.addWidget(self.btn_logar)
        layout_login.addStretch()
        
        # ABA CADASTRO
        self.tab_cadastro = QWidget()
        layout_cad = QVBoxLayout(self.tab_cadastro)
        
        import idiomas
        self.chk_gringo = QCheckBox(idiomas.tr("cad_chk_gringo") if hasattr(idiomas, 'tr') else "Sou de fora do Brasil / I'm an international user")
        self.chk_gringo.setStyleSheet("color: #ffc107; font-weight: bold; margin-bottom: 5px;")
        self.chk_gringo.toggled.connect(self.alternar_modo_gringo)

        self.inp_cad_nome = QLineEdit()
        self.inp_cad_nome.setPlaceholderText("Nome Completo")
        self.inp_cad_email = QLineEdit()
        self.inp_cad_email.setPlaceholderText("E-mail")
        
        # Container de Documento (Brasil) vs País (Gringo)
        self.container_doc = QWidget()
        lay_doc = QVBoxLayout(self.container_doc)
        lay_doc.setContentsMargins(0, 0, 0, 0)
        self.lbl_doc = QLabel("CPF/CNPJ:")
        self.inp_cad_doc = QLineEdit()
        self.inp_cad_doc.setPlaceholderText("Apenas números")
        lay_doc.addWidget(self.lbl_doc)
        lay_doc.addWidget(self.inp_cad_doc)

        self.container_pais = QWidget()
        lay_pais = QVBoxLayout(self.container_pais)
        lay_pais.setContentsMargins(0, 0, 0, 0)
        self.lbl_pais = QLabel("Country:")
        self.combo_pais = QComboBox()
        self.combo_pais.addItems(["United States", "Canada", "United Kingdom", "Portugal", "Australia", "Other"])
        lay_pais.addWidget(self.lbl_pais)
        lay_pais.addWidget(self.combo_pais)
        self.container_pais.hide() # Escondido por padrão

        self.inp_cad_senha = QLineEdit()
        self.inp_cad_senha.setPlaceholderText("Crie uma Senha")
        self.inp_cad_senha.setEchoMode(QLineEdit.Password)
        self.btn_cadastrar = QPushButton("Criar Conta")
        self.btn_cadastrar.setStyleSheet("background-color: #28a745; color: white; padding: 10px; font-weight: bold;")
        self.btn_cadastrar.clicked.connect(self.fazer_cadastro)
        
        layout_cad.addWidget(QLabel("Novo por aqui? Cadastre-se:"))
        layout_cad.addWidget(self.chk_gringo)
        layout_cad.addWidget(QLabel("Nome Completo:"))
        layout_cad.addWidget(self.inp_cad_nome)
        layout_cad.addWidget(QLabel("E-mail:"))
        layout_cad.addWidget(self.inp_cad_email)
        layout_cad.addWidget(self.container_doc)
        layout_cad.addWidget(self.container_pais)
        layout_cad.addWidget(QLabel("Senha:"))
        layout_cad.addWidget(self.inp_cad_senha)
        layout_cad.addWidget(self.btn_cadastrar)
        layout_cad.addStretch()

        self.tabs.addTab(self.tab_login, "Login")
        self.tabs.addTab(self.tab_cadastro, "Criar Conta")
        layout.addWidget(self.tabs)

    def alternar_modo_gringo(self, checked):
        if checked:
            self.container_doc.hide()
            self.container_pais.show()
        else:
            self.container_pais.hide()
            self.container_doc.show()

    def fazer_login(self):
        email = self.inp_login_email.text().strip()
        senha = self.inp_login_senha.text()
        if not email or not senha:
            QMessageBox.warning(self, "Aviso", "Preencha e-mail e senha.")
            return
        self.btn_logar.setEnabled(False)
        self.btn_logar.setText("Autenticando...")
        
        payload = {"email": email, "senha": senha}
        self.worker = WorkerAuthApp('login', payload)
        self.worker.sucesso.connect(self.auth_sucesso)
        self.worker.erro.connect(self.auth_erro)
        self.worker.start()

    def fazer_cadastro(self):
        nome = self.inp_cad_nome.text().strip()
        email = self.inp_cad_email.text().strip()
        senha = self.inp_cad_senha.text()
        
        is_gringo = self.chk_gringo.isChecked()
        doc = "".join(filter(str.isdigit, self.inp_cad_doc.text())) if not is_gringo else "EXTRANGEIRO"
        pais = self.combo_pais.currentText() if is_gringo else "BR"
        
        if not nome or not email or not senha or (not is_gringo and not doc):
            QMessageBox.warning(self, "Aviso", "Preencha todos os campos obrigatórios.")
            return
            
        self.btn_cadastrar.setEnabled(False)
        self.btn_cadastrar.setText("Registrando...")
        
        payload = {
            "nome": nome, 
            "email": email, 
            "documento": doc, 
            "senha": senha, 
            "produto": "SmartSlides Pro",
            "pais": pais
        }
        self.worker = WorkerAuthApp('registrar', payload)
        self.worker.sucesso.connect(self.auth_sucesso)
        self.worker.erro.connect(self.auth_erro)
        self.worker.start()

    def auth_sucesso(self, token):
        self.token_recebido = token
        QMessageBox.information(self, "Sucesso", "Conta conectada com sucesso!")
        self.accept()

    def auth_erro(self, msg):
        self.btn_logar.setEnabled(True)
        self.btn_logar.setText("Entrar")
        self.btn_cadastrar.setEnabled(True)
        self.btn_cadastrar.setText("Criar Conta")
        QMessageBox.critical(self, "Aviso", msg)

# =========================================================
# LOJA UNIVERSAL (VITRINE, CHECKOUT E CUPÕES)
# =========================================================
class WorkerVitrine(QThread):
    sucesso = Signal(dict)
    erro = Signal(str)
    
    def __init__(self, moeda="BRL"):
        super().__init__()
        self.moeda = moeda

    def run(self):
        try:
            # Pede a vitrine já informando se quer os preços do Stripe ou do Asaas
            url = f"https://vegap-masterapp.hf.space/master/vitrine/SmartSlides%20Pro?moeda={self.moeda}"
            resp = requests.get(url)
            if resp.status_code == 200:
                self.sucesso.emit(resp.json())
            else:
                self.erro.emit("Erro ao carregar a vitrine.")
        except Exception as e:
            self.erro.emit(str(e))

class WorkerCheckout(QThread):
    sucesso = Signal(dict)
    erro = Signal(str)

    def __init__(self, token, produto_id, plano_nome, cupom, gateway="asaas"):
        super().__init__()
        self.payload = {
            "token": token,
            "produto_id": produto_id,
            "plano_nome": plano_nome,
            "cupom": cupom,
            "gateway": gateway # Diz pro backend se é pra gerar link do Asaas ou do Stripe
        }

    def run(self):
        try:
            resp = requests.post("https://vegap-masterapp.hf.space/master/pagamentos/gerar-checkout", json=self.payload)
            if resp.status_code == 200:
                self.sucesso.emit(resp.json())
            else:
                self.erro.emit(resp.json().get("detail", "Erro desconhecido"))
        except Exception as e:
            self.erro.emit(str(e))

class WorkerStatusPagamento(QThread):
    sucesso = Signal()

    def __init__(self, cobranca_id):
        super().__init__()
        self.cobranca_id = cobranca_id
        self.rodando = True

    def run(self):
        while self.rodando:
            try:
                resp = requests.get(f"https://vegap-masterapp.hf.space/master/pagamentos/checar-status?cobranca_id={self.cobranca_id}")
                if resp.status_code == 200 and resp.json().get("pago"):
                    self.sucesso.emit()
                    break
            except:
                pass
            time.sleep(5)

    def stop(self):
        self.rodando = False

class LojaDialog(QDialog):
    def __init__(self, token_cliente, filtro_loja="SAAS", parent=None):
        super().__init__(parent)
        self.setWindowTitle("🛒 Loja VegaTech - SmartSlides Pro")
        self.setFixedSize(450, 500)
        self.token_cliente = token_cliente
        self.filtro_loja = filtro_loja # "SAAS" ou "BYOK"
        self.produto_id = None
        self.worker_status = None
        self.plano_em_andamento = None
        self.plano_comprado = None # Guarda o recibo final
        self.init_ui()
        self.carregar_vitrine()

    def init_ui(self):
        self.layout = QVBoxLayout(self)
        
        import idiomas
        self.lbl_titulo = QLabel(idiomas.tr("loja_loading"))
        self.lbl_titulo.setAlignment(Qt.AlignCenter)
        self.lbl_titulo.setStyleSheet("font-size: 18px; font-weight: bold; margin-bottom: 10px;")
        self.layout.addWidget(self.lbl_titulo)
        
        # CHAVINHA DE MOEDA (BRL / USD)
        hbox_moeda = QHBoxLayout()
        self.lbl_moeda = QLabel("Moeda (Currency):")
        self.combo_moeda = QComboBox()
        self.combo_moeda.addItems(["BRL (R$) - Brasil", "USD ($) - International"])
        self.combo_moeda.setStyleSheet("padding: 4px; font-weight: bold;")
        self.combo_moeda.currentIndexChanged.connect(self.carregar_vitrine)
        hbox_moeda.addWidget(self.lbl_moeda)
        hbox_moeda.addWidget(self.combo_moeda)
        hbox_moeda.addStretch()
        self.layout.addLayout(hbox_moeda)
        
        # --- NOVO: SCROLL AREA PARA OS PLANOS NÃO ESMAGAREM ---
        self.scroll_planos = QScrollArea()
        self.scroll_planos.setWidgetResizable(True)
        self.scroll_planos.setStyleSheet("QScrollArea { border: none; background-color: transparent; }")
        
        self.widget_planos = QWidget()
        self.widget_planos.setStyleSheet("background-color: transparent;")
        self.container_planos = QVBoxLayout(self.widget_planos)
        self.container_planos.setContentsMargins(0, 0, 10, 0) # Margem pra barra de rolagem não colar
        self.container_planos.setSpacing(10)
        
        self.scroll_planos.setWidget(self.widget_planos)
        self.layout.addWidget(self.scroll_planos)
        # ------------------------------------------------------
        
        self.inp_cupom = QLineEdit()
        self.inp_cupom.setPlaceholderText(idiomas.tr("loja_cupom"))
        self.inp_cupom.setStyleSheet("font-size: 14px; padding: 10px; border-radius: 5px; border: 1px solid gray;")
        self.layout.addWidget(self.inp_cupom)
        
        self.lbl_status = QLabel("")
        self.lbl_status.setAlignment(Qt.AlignCenter)
        self.lbl_status.setStyleSheet("color: #0078d7; font-weight: bold; margin-top: 10px;")
        self.layout.addWidget(self.lbl_status)

    def carregar_vitrine(self):
        import idiomas
        self.lbl_titulo.setText(idiomas.tr("loja_loading"))
        
        while self.container_planos.count():
            item = self.container_planos.takeAt(0)
            if item.widget(): item.widget().deleteLater()
            
        moeda_escolhida = "USD" if self.combo_moeda.currentIndex() == 1 else "BRL"
        
        self.worker_vitrine = WorkerVitrine(moeda=moeda_escolhida)
        self.worker_vitrine.sucesso.connect(self.montar_planos)
        self.worker_vitrine.erro.connect(lambda e: self.lbl_titulo.setText("Erro ao carregar a loja."))
        self.worker_vitrine.start()

    def montar_planos(self, dados):
        from PySide6.QtWidgets import QFrame # Garante a importação do QFrame para o visual
        
        self.produto_id = dados.get("id")
        planos = dados.get("planos_precos", {})
        
        # Limpa o container caso a pessoa clique duas vezes
        while self.container_planos.count():
            item = self.container_planos.takeAt(0)
            if item.widget(): item.widget().deleteLater()
        
        import idiomas
        if not planos:
            self.lbl_titulo.setText(idiomas.tr("loja_vazio"))
            return
            
        self.lbl_titulo.setText(idiomas.tr("loja_escolha"))
        
        # MÁGICA DA ORDENAÇÃO: Força a vitrine a seguir a hierarquia clássica de SaaS
        ordem_ciclos = {"mensal": 1, "trimestral": 2, "semestral": 3, "anual": 4, "unico": 5}
        planos_ordenados = sorted(
            planos.items(), 
            key=lambda x: ordem_ciclos.get(str(x[1].get("ciclo")).lower(), 99)
        )
        
        for nome_plano, detalhes in planos_ordenados:
            # GOLPE VEGATECH: Filtro Absoluto! 
            if self.filtro_loja == "BYOK" and "BYOK" not in nome_plano.upper():
                continue
            # Se pediu SAAS genérico, não mostra o BYOK
            if self.filtro_loja == "SAAS" and "BYOK" in nome_plano.upper():
                continue
            # A VITRINE VIP: Se pediu ELITE, esconde tudo que não tiver ELITE no nome
            if self.filtro_loja == "ELITE" and "ELITE" not in nome_plano.upper():
                continue
                
            valor = float(detalhes.get("valor", 0.0))
            trial_dias = int(detalhes.get("trial_dias", 0))
            ciclo = detalhes.get("ciclo", "unico").capitalize()
            
            # MÁGICA VISUAL: Corta o "_MENSAL" do nome para o cliente ler só "SAAS"
            nome_exibicao = nome_plano.split("_")[0] if "_" in nome_plano else nome_plano
            
            # CRIA UM CARD CHIQUE DE ALTA CONVERSÃO
            card = QFrame()
            card.setStyleSheet("background-color: #262626; border: 1px solid #444; border-radius: 12px; margin-bottom: 5px;")
            card_layout = QVBoxLayout(card)
            
            # Título e Preço (Usa o nome maquiado)
            lbl_nome = QLabel(f"⭐ {nome_exibicao.upper()}")
            lbl_nome.setStyleSheet("font-size: 16px; font-weight: 900; color: white; border: none;")
            
            txt_ciclo = f" / {ciclo}" if ciclo != "Unico" else ""
            simbolo_moeda = "$" if self.combo_moeda.currentIndex() == 1 else "R$"
            lbl_preco = QLabel(f"{simbolo_moeda} {valor:.2f}{txt_ciclo}")
            lbl_preco.setStyleSheet("font-size: 24px; font-weight: 800; color: #28a745; border: none;")
            
            card_layout.addWidget(lbl_nome)
            card_layout.addWidget(lbl_preco)
            
            # GATILHO MENTAL DO TRIAL (O SEGREDO DA CONVERSÃO!)
            if trial_dias > 0:
                lbl_trial = QLabel(f"🎁 {trial_dias} DIAS TOTALMENTE GRÁTIS!")
                lbl_trial.setStyleSheet("""
                    font-size: 13px; font-weight: bold; color: #ffc107; 
                    background-color: rgba(255, 193, 7, 0.1); 
                    padding: 6px; border-radius: 4px; border: 1px solid #ffc107;
                """)
                lbl_trial.setAlignment(Qt.AlignCenter)
                card_layout.addWidget(lbl_trial)
            
            # Botão de Assinar
            btn_comprar = QPushButton(idiomas.tr("loja_btn_assinar"))
            btn_comprar.setMinimumHeight(45) # <--- BLINDAGEM CONTRA ESMAGAMENTO
            btn_comprar.setStyleSheet("background-color: #0078d7; color: white; font-weight: bold; font-size: 15px; border-radius: 6px; border: none;")
            btn_comprar.setCursor(Qt.PointingHandCursor)
            btn_comprar.clicked.connect(lambda checked=False, p=nome_plano: self.comprar_plano(p))
            
            card_layout.addWidget(btn_comprar)
            self.container_planos.addWidget(card)

    def comprar_plano(self, plano_nome):
        if not self.token_cliente:
            dialog_auth = LoginCadastroDialog(self)
            if dialog_auth.exec() == QDialog.Accepted:
                self.token_cliente = dialog_auth.token_recebido
            else:
                return

        import idiomas
        cupom = self.inp_cupom.text().strip()
        self.lbl_status.setText(idiomas.tr("loja_gerando"))
        self.setEnabled(False)
        
        gateway_escolhido = "stripe" if self.combo_moeda.currentIndex() == 1 else "asaas"
        
        self.plano_em_andamento = plano_nome # Anota qual botão ele clicou
        self.worker_checkout = WorkerCheckout(self.token_cliente, self.produto_id, plano_nome, cupom, gateway_escolhido)
        self.worker_checkout.sucesso.connect(self.processar_checkout)
        self.worker_checkout.erro.connect(self.erro_checkout)
        self.worker_checkout.start()

    def processar_checkout(self, dados):
        if dados.get("gratuito"):
            self.plano_comprado = self.plano_em_andamento # Grava o recibo!
            QMessageBox.information(self, "Sucesso", dados.get("mensagem", "Acesso VIP liberado com sucesso!"))
            self.accept()
            return
            
        checkout_url = dados.get("checkout_url")
        cobranca_id = dados.get("cobranca_id")
        
        if checkout_url:
            webbrowser.open(checkout_url)
            self.lbl_status.setText("Navegador aberto! A aguardar o pagamento...")
            self.setEnabled(True)
            
            self.worker_status = WorkerStatusPagamento(cobranca_id)
            self.worker_status.sucesso.connect(self.pagamento_confirmado)
            self.worker_status.start()

    def erro_checkout(self, erro):
        self.setEnabled(True)
        self.lbl_status.setText("")
        QMessageBox.critical(self, "Erro", erro)

    def pagamento_confirmado(self):
        self.plano_comprado = self.plano_em_andamento # Grava o recibo!
        QMessageBox.information(self, "Sucesso", "Pagamento confirmado! Acesso Liberado!")
        self.accept()

    def closeEvent(self, event):
        if self.worker_status:
            self.worker_status.stop()
        super().closeEvent(event)

# --- JANELA DE CONFIGURAÇÕES ---
class ConfigDialog(QDialog):
    def __init__(self, current_config, parent=None):
        super().__init__(parent)
        self.setWindowTitle("⚙️ Configurações do App")
        self.resize(500, 400)
        self.config = current_config
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout(self)
        self.tabs = QTabWidget()
        
        # Aba 1: Licença
        self.tab_licenca = QWidget()
        vbox_lic = QVBoxLayout(self.tab_licenca)
        
        # O TÍTULO EXPLÍCITO DA ABA
        self.lbl_explicacao_licenca = QLabel()
        self.lbl_explicacao_licenca.setWordWrap(True)
        self.lbl_explicacao_licenca.setStyleSheet("margin-bottom: 10px;") 
        vbox_lic.addWidget(self.lbl_explicacao_licenca)
        
        self.lbl_status_lic = QLabel()
        vbox_lic.addWidget(self.lbl_status_lic)
        self.combo_licenca = QComboBox()
        self.combo_licenca.addItem("FREE (Edição Manual)", "FREE (Edição Manual)")
        self.combo_licenca.addItem("BYOK (Sua Chave)", "BYOK (Sua Chave)")
        self.combo_licenca.addItem("SAAS (Chave Embutida)", "SAAS (Chave Embutida)")
        
        lic_atual = self.config.get("license", "FREE")
        if "SAAS" in lic_atual: lic_atual = "SAAS (Chave Embutida)" # Conserta se ficou salvo o lixo da versão anterior
        
        self.combo_licenca.setCurrentText(lic_atual)
        self.combo_licenca.currentTextChanged.connect(self.update_ia_tab_visibility)
        vbox_lic.addWidget(self.combo_licenca)
        vbox_lic.addStretch()
        
        # Aba 2: Inteligência Artificial
        self.tab_ia = QWidget()
        self.vbox_ia = QVBoxLayout(self.tab_ia)
        
        self.lbl_aviso_saas = QLabel("✅ <b>Licença SaaS Ativa</b>\nO motor principal está liberado.")
        self.lbl_aviso_saas.setStyleSheet("color: #28a745; margin-bottom: 5px;")
        self.vbox_ia.addWidget(self.lbl_aviso_saas)

        # --- OPÇÃO 1: GOOGLE GEMINI ---
        self.chk_gemini = QCheckBox("🔵 Usar Google Gemini (Padrão)")
        self.chk_gemini.setStyleSheet("font-weight: bold; font-size: 13px;")
        self.vbox_ia.addWidget(self.chk_gemini)
        
        self.container_gemini = QWidget()
        vbox_gem = QVBoxLayout(self.container_gemini)
        vbox_gem.setContentsMargins(15, 0, 0, 10)
        
        self.wrapper_keys_gemini = QWidget()
        lay_keys_gem = QVBoxLayout(self.wrapper_keys_gemini)
        lay_keys_gem.setContentsMargins(0, 0, 0, 0)
        self.lbl_key_gemini = QLabel("Chave de API Gemini (BYOK):")
        lay_keys_gem.addWidget(self.lbl_key_gemini)
        self.txt_api_key = QLineEdit(self.config.get("api_key", ""))
        self.txt_api_key.setEchoMode(QLineEdit.Password)
        lay_keys_gem.addWidget(self.txt_api_key)
        vbox_gem.addWidget(self.wrapper_keys_gemini)
        
        hbox_mod_gem = QHBoxLayout()
        self.combo_model = QComboBox()
        self.combo_model.addItem(self.config.get("model", "gemini-2.5-flash"))
        self.btn_load_models = QPushButton("🔄 Buscar Modelos")
        self.btn_load_models.clicked.connect(self.load_gemini_models_live)
        self.lbl_mod_gem = QLabel("Modelo:")
        hbox_mod_gem.addWidget(self.lbl_mod_gem)
        hbox_mod_gem.addWidget(self.combo_model, 1)
        hbox_mod_gem.addWidget(self.btn_load_models)
        vbox_gem.addLayout(hbox_mod_gem)
        self.vbox_ia.addWidget(self.container_gemini)

        # --- OPÇÃO 2: CLOUDFLARE WORKERS AI ---
        hbox_cf_title = QHBoxLayout()
        self.chk_cloudflare = QCheckBox("🟠 Usar Cloudflare Workers AI")
        self.chk_cloudflare.setStyleSheet("font-weight: bold; font-size: 13px;")
        hbox_cf_title.addWidget(self.chk_cloudflare)
        
        self.btn_upgrade_elite = QPushButton("⭐ Upgrade ELITE")
        self.btn_upgrade_elite.setStyleSheet("background-color: #ffc107; color: black; font-weight: bold; padding: 2px 8px; border-radius: 4px;")
        self.btn_upgrade_elite.setCursor(Qt.PointingHandCursor)
        self.btn_upgrade_elite.clicked.connect(self.abrir_loja_upgrade)
        self.btn_upgrade_elite.hide()
        
        hbox_cf_title.addWidget(self.btn_upgrade_elite)
        hbox_cf_title.addStretch()
        self.vbox_ia.addLayout(hbox_cf_title)
        
        self.container_cf = QWidget()
        vbox_cf = QVBoxLayout(self.container_cf)
        vbox_cf.setContentsMargins(15, 0, 0, 10)
        
        self.wrapper_keys_cf = QWidget()
        lay_keys_cf = QVBoxLayout(self.wrapper_keys_cf)
        lay_keys_cf.setContentsMargins(0, 0, 0, 0)
        self.lbl_acc_cf = QLabel("Cloudflare Account ID (BYOK):")
        lay_keys_cf.addWidget(self.lbl_acc_cf)
        self.txt_cf_account = QLineEdit(self.config.get("cf_account_id", ""))
        lay_keys_cf.addWidget(self.txt_cf_account)
        self.lbl_tok_cf = QLabel("Cloudflare API Token (BYOK):")
        lay_keys_cf.addWidget(self.lbl_tok_cf)
        self.txt_cf_token = QLineEdit(self.config.get("cf_api_token", ""))
        self.txt_cf_token.setEchoMode(QLineEdit.Password)
        lay_keys_cf.addWidget(self.txt_cf_token)
        vbox_cf.addWidget(self.wrapper_keys_cf)
        
        hbox_mod_cf = QHBoxLayout()
        self.combo_cf_model = QComboBox()
        self.combo_cf_model.addItem(self.config.get("cf_model", "@cf/meta/llama-3-8b-instruct"))
        self.btn_load_cf = QPushButton("🔄 Carregar CF")
        self.btn_load_cf.clicked.connect(self.load_cf_models_live)
        self.lbl_mod_cf = QLabel("Modelo:")
        hbox_mod_cf.addWidget(self.lbl_mod_cf)
        hbox_mod_cf.addWidget(self.combo_cf_model, 1)
        hbox_mod_cf.addWidget(self.btn_load_cf)
        vbox_cf.addLayout(hbox_mod_cf)
        self.vbox_ia.addWidget(self.container_cf)
        self.vbox_ia.addStretch()

        # --- Aba 3: Aparência ---
        self.tab_visual = QWidget()
        vbox_vis = QVBoxLayout(self.tab_visual)
        
        self.lbl_idioma = QLabel()
        vbox_vis.addWidget(self.lbl_idioma)
        self.combo_idioma = QComboBox()
        self.combo_idioma.addItem("Português", "pt")
        self.combo_idioma.addItem("English", "en")
        idx_idioma = self.combo_idioma.findData(self.config.get("idioma", "pt"))
        if idx_idioma >= 0: self.combo_idioma.setCurrentIndex(idx_idioma)
        self.combo_idioma.currentIndexChanged.connect(self.mudar_idioma_ao_vivo)
        vbox_vis.addWidget(self.combo_idioma)

        self.lbl_tema = QLabel()
        self.lbl_tema.setStyleSheet("margin-top: 10px;")
        vbox_vis.addWidget(self.lbl_tema)
        self.combo_theme = QComboBox()
        self.combo_theme.addItem("Escuro", "Escuro")
        self.combo_theme.addItem("Claro", "Claro")
        idx_tema = self.combo_theme.findData(self.config.get("theme", "Escuro"))
        if idx_tema >= 0: self.combo_theme.setCurrentIndex(idx_tema)
        vbox_vis.addWidget(self.combo_theme)
        
        self.lbl_zoom = QLabel()
        self.lbl_zoom.setStyleSheet("margin-top: 10px;")
        vbox_vis.addWidget(self.lbl_zoom)
        self.slider_zoom = QSlider(Qt.Horizontal)
        self.slider_zoom.setRange(0, 150)
        self.slider_zoom.setValue(self.config.get("zoom", 100))
        self.lbl_zoom_val = QLabel(f"{self.slider_zoom.value()}%")
        self.slider_zoom.valueChanged.connect(self.aplicar_zoom_ao_vivo)
        
        hbox_zoom = QHBoxLayout()
        hbox_zoom.addWidget(self.slider_zoom)
        hbox_zoom.addWidget(self.lbl_zoom_val)
        vbox_vis.addLayout(hbox_zoom)
        vbox_vis.addStretch()

        self.tabs.addTab(self.tab_licenca, "Licença")
        self.tabs.addTab(self.tab_ia, "Inteligência Artificial")
        self.tabs.addTab(self.tab_visual, "Aparência")
        layout.addWidget(self.tabs)

        self.btn_salvar = QPushButton()
        self.btn_salvar.setStyleSheet("background-color: #0078d7; color: white; padding: 10px; font-weight: bold;")
        self.btn_salvar.clicked.connect(self.save_and_close)
        layout.addWidget(self.btn_salvar)

        self.chk_gemini.toggled.connect(lambda checked: self.alternar_provedor("gemini", checked))
        self.chk_cloudflare.toggled.connect(lambda checked: self.alternar_provedor("cloudflare", checked))

        prov_atual = self.config.get("provedor", "gemini")
        if prov_atual == "cloudflare":
            self.chk_cloudflare.setChecked(True)
        else:
            self.chk_gemini.setChecked(True)
            
        self.update_ia_tab_visibility(self.combo_licenca.currentText())
        self.atualizar_textos_ui()
        
    def mudar_idioma_ao_vivo(self):
        import idiomas
        idiomas.set_idioma(self.combo_idioma.currentData())
        self.atualizar_textos_ui()
        if self.parent() and hasattr(self.parent(), 'atualizar_textos_ui'):
            self.parent().atualizar_textos_ui()

    def atualizar_textos_ui(self):
        import idiomas
        self.setWindowTitle(idiomas.tr("cfg_title"))
        self.lbl_explicacao_licenca.setText(idiomas.tr("cfg_lbl_lic"))
        self.lbl_status_lic.setText(idiomas.tr("cfg_status_lic"))
        self.lbl_idioma.setText(idiomas.tr("cfg_lbl_idioma"))
        self.lbl_tema.setText(idiomas.tr("cfg_lbl_tema"))
        self.lbl_zoom.setText(idiomas.tr("cfg_lbl_zoom"))
        self.btn_salvar.setText(idiomas.tr("cfg_btn_salvar"))
        self.tabs.setTabText(0, idiomas.tr("cfg_tab_lic"))
        self.tabs.setTabText(1, idiomas.tr("cfg_tab_ia"))
        self.tabs.setTabText(2, idiomas.tr("cfg_tab_vis"))
        self.btn_load_models.setText(idiomas.tr("cfg_btn_buscar"))
        self.btn_load_cf.setText(idiomas.tr("cfg_btn_load_cf"))
        self.combo_licenca.setItemText(0, idiomas.tr("lic_free"))
        self.combo_licenca.setItemText(1, idiomas.tr("lic_byok"))
        self.combo_licenca.setItemText(2, idiomas.tr("lic_saas"))
        
        self.combo_theme.setItemText(0, idiomas.tr("tema_escuro"))
        self.combo_theme.setItemText(1, idiomas.tr("tema_claro"))

        self.lbl_aviso_saas.setText(idiomas.tr("cfg_aviso_saas"))
        self.chk_gemini.setText(idiomas.tr("cfg_chk_gemini"))
        self.chk_cloudflare.setText(idiomas.tr("cfg_chk_cf"))
        self.btn_upgrade_elite.setText(idiomas.tr("cfg_btn_upgrade"))
        self.lbl_key_gemini.setText(idiomas.tr("cfg_lbl_key_gemini"))
        self.lbl_mod_gem.setText(idiomas.tr("cfg_lbl_modelo"))
        self.lbl_acc_cf.setText(idiomas.tr("cfg_lbl_acc_cf"))
        self.lbl_tok_cf.setText(idiomas.tr("cfg_lbl_tok_cf"))
        self.lbl_mod_cf.setText(idiomas.tr("cfg_lbl_modelo"))

    def aplicar_zoom_ao_vivo(self, v):
        self.lbl_zoom_val.setText(f"{v}%")
        # Puxa o cordão da janela Pai e aplica o zoom na hora que deslizar!
        if self.parent() and hasattr(self.parent(), 'apply_zoom'):
            self.parent().apply_zoom(v)

    def update_ia_tab_visibility(self, licensa):
        if "SAAS" in licensa:
            self.lbl_aviso_saas.show()
            self.wrapper_keys_gemini.hide()
            self.wrapper_keys_cf.hide()
            self.tabs.setTabEnabled(1, True)
            
            if not self.config.get("saas_elite_ativo"):
                self.chk_cloudflare.setEnabled(False)
                self.btn_upgrade_elite.show()
                if hasattr(self, 'chk_cloudflare') and self.chk_cloudflare.isChecked():
                    self.chk_gemini.setChecked(True)
            else:
                self.chk_cloudflare.setEnabled(True)
                self.btn_upgrade_elite.hide()
                
        elif "BYOK" in licensa:
            self.lbl_aviso_saas.hide()
            self.wrapper_keys_gemini.show()
            self.wrapper_keys_cf.show()
            self.tabs.setTabEnabled(1, True)
            if hasattr(self, 'chk_cloudflare'):
                self.chk_cloudflare.setEnabled(True)
                self.btn_upgrade_elite.hide()
        else:
            self.tabs.setTabEnabled(1, False) 

    def abrir_loja_upgrade(self):
        if not self.config.get("token_master"):
            QMessageBox.warning(self, "Aviso", "Autenticação necessária. Salve as configurações primeiro para fazer login.")
            return
            
        # AQUI ESTÁ A CORREÇÃO: Manda abrir a loja passando "ELITE" no filtro
        loja = LojaDialog(self.config.get("token_master"), "ELITE", self)
        if loja.exec() == QDialog.Accepted:
            self.config["token_master"] = loja.token_cliente
            if loja.plano_comprado and "ELITE" in loja.plano_comprado.upper():
                self.config["saas_elite_ativo"] = True
                self.config["saas_padrao_ativo"] = False
                QMessageBox.information(self, "Bem-vindo ao Elite!", "Poder Multi-IA liberado com sucesso!")
                self.update_ia_tab_visibility("SAAS (Chave Embutida)")

    def alternar_provedor(self, origem, checked):
        if not checked:
            if origem == "gemini" and not self.chk_cloudflare.isChecked():
                self.chk_gemini.blockSignals(True)
                self.chk_gemini.setChecked(True)
                self.chk_gemini.blockSignals(False)
            elif origem == "cloudflare" and not self.chk_gemini.isChecked():
                self.chk_cloudflare.blockSignals(True)
                self.chk_cloudflare.setChecked(True)
                self.chk_cloudflare.blockSignals(False)
            return

        if origem == "gemini":
            self.chk_cloudflare.blockSignals(True)
            self.chk_cloudflare.setChecked(False)
            self.chk_cloudflare.blockSignals(False)
            self.container_gemini.setEnabled(True)
            self.container_cf.setEnabled(False) 
        else:
            self.chk_gemini.blockSignals(True)
            self.chk_gemini.setChecked(False)
            self.chk_gemini.blockSignals(False)
            self.container_cf.setEnabled(True)
            self.container_gemini.setEnabled(False)

    def load_gemini_models_live(self):
        licensa = self.combo_licenca.currentText()
        chave = self.txt_api_key.text().strip() if licensa == "BYOK (Sua Chave)" else None
        
        if licensa == "BYOK (Sua Chave)" and not chave:
            QMessageBox.warning(self, "Aviso", "Digite a chave BYOK para consultar modelos.")
            return
            
        self.btn_load_models.setText("⏳ Buscando...")
        QApplication.processEvents()
        
        modelos, erro = motor_ia.listar_modelos_gemini(chave)
        self.combo_model.clear()
        if modelos: 
            self.combo_model.addItems(modelos)
        else: 
            QMessageBox.critical(self, "Erro na API Gemini", f"A API do Google recusou a conexão.\n\nDetalhe do erro:\n{erro}")
        self.btn_load_models.setText("🔄 Buscar Modelos")

    def load_cf_models_live(self):
        acc = self.txt_cf_account.text().strip()
        tok = self.txt_cf_token.text().strip()
        self.btn_load_cf.setText("⏳ Carregando...")
        QApplication.processEvents()
        
        modelos, status = motor_ia.listar_modelos_cloudflare(acc, tok)
        self.combo_cf_model.clear()
        
        if modelos:
            self.combo_cf_model.addItems(modelos)
        else:
            QMessageBox.critical(self, "Erro na API Cloudflare", f"Não foi possível carregar os modelos.\n\nDetalhe:\n{status}")
            
        self.btn_load_cf.setText("🔄 Carregar CF")

    def save_and_close(self):
        licenca_escolhida = self.combo_licenca.currentData()
        licenca_salva = self.config.get("license", "FREE")
        
        # 1. Barreira do BYOK (Com Reversão Automática Anti-Bypass)
        if licenca_escolhida == "BYOK (Sua Chave)":
            if not self.config.get("byok_ativo"): 
                if not self.config.get("token_master"):
                    dialog_auth = LoginCadastroDialog(self)
                    if dialog_auth.exec() == QDialog.Accepted:
                        self.config["token_master"] = dialog_auth.token_recebido
                    else: 
                        self.combo_licenca.setCurrentText(licenca_salva)
                        return
                
                loja = LojaDialog(self.config.get("token_master"), "BYOK", self)
                if loja.exec() == QDialog.Accepted:
                    self.config["token_master"] = loja.token_cliente
                    self.config["byok_ativo"] = True
                else: 
                    self.combo_licenca.setCurrentText(licenca_salva)
                    return

        # 2. Barreira do SAAS (Com Reversão Automática Anti-Bypass)
        elif licenca_escolhida == "SAAS (Chave Embutida)":
            if not self.config.get("saas_padrao_ativo") and not self.config.get("saas_elite_ativo"):
                if not self.config.get("token_master"):
                    dialog_auth = LoginCadastroDialog(self)
                    if dialog_auth.exec() == QDialog.Accepted:
                        self.config["token_master"] = dialog_auth.token_recebido
                    else: 
                        self.combo_licenca.setCurrentText(licenca_salva)
                        return
                
                loja = LojaDialog(self.config.get("token_master"), "SAAS", self)
                if loja.exec() == QDialog.Accepted:
                    self.config["token_master"] = loja.token_cliente
                    if loja.plano_comprado and "ELITE" in loja.plano_comprado.upper():
                        self.config["saas_elite_ativo"] = True
                        self.config["saas_padrao_ativo"] = False
                    else:
                        self.config["saas_padrao_ativo"] = True
                        self.config["saas_elite_ativo"] = False
                else: 
                    self.combo_licenca.setCurrentText(licenca_salva)
                    return

        # Se passou e pagou (ou já tinha pago), salva limpo
        self.config["license"] = licenca_escolhida
        self.config["theme"] = self.combo_theme.currentData()
        self.config["zoom"] = self.slider_zoom.value()
        self.config["idioma"] = self.combo_idioma.currentData()
        
        self.config["provedor"] = "gemini" if self.chk_gemini.isChecked() else "cloudflare"
        self.config["api_key"] = self.txt_api_key.text().strip()
        self.config["model"] = self.combo_model.currentText()
        
        self.config["cf_account_id"] = self.txt_cf_account.text().strip()
        self.config["cf_api_token"] = self.txt_cf_token.text().strip()
        self.config["cf_model"] = self.combo_cf_model.currentText()
        
        # Garante a gravação permanente das chaves de ativação na memória
        self.config["saas_elite_ativo"] = self.config.get("saas_elite_ativo", False)
        self.config["saas_padrao_ativo"] = self.config.get("saas_padrao_ativo", False)
        self.config["byok_ativo"] = self.config.get("byok_ativo", False)

        self.accept()        
class WorkerIAGerador(QThread):
    progresso = Signal(int, str)
    concluido = Signal()

    def __init__(self, tema, config_ia, roteiro_carregado=None):
        super().__init__()
        self.tema = tema
        self.config_ia = config_ia
        self.roteiro_carregado = roteiro_carregado
        self.roteiro = None
        self.slide_images_dict = {}
        self.erro_msg = None

    def run(self):
        try:
            if self.roteiro_carregado:
                self.roteiro = self.roteiro_carregado
            else:
                self.progresso.emit(5, "Conectando ao motor de IA...")
                sucesso, roteiro_ou_erro = motor_ia.gerar_roteiro_slides(self.tema, self.config_ia)
                
                if self.isInterruptionRequested(): return
                
                if not sucesso:
                    self.erro_msg = roteiro_ou_erro
                    return
                    
                self.roteiro = roteiro_ou_erro

            # 2. Baixando as Imagens
            total_slides = len(self.roteiro)
            self.progresso.emit(20, "Estrutura pronta! Buscando imagens...")
            
            for i, slide_data in enumerate(self.roteiro):
                if self.isInterruptionRequested(): return
                palavra_chave = slide_data.get("palavra_chave_imagem", "")
                if palavra_chave:
                    self.progresso.emit(20 + int(((i + 1) / total_slides) * 65), f"Baixando opções para o slide {i + 1} de {total_slides}...")
                    paths = motor_ia.baixar_imagem_pexels(palavra_chave, i)
                    if paths:
                        self.slide_images_dict[i] = paths

            self.progresso.emit(90, "Abrindo painel de escolhas...")
        except Exception as e:
            import traceback
            self.erro_msg = f"Erro no Worker:\n{str(e)}\n\n{traceback.format_exc()}"
        finally:
            self.concluido.emit()
            
# --- JANELA COM ROLAGEM PARA NOVIDADES ---
class NovidadesDialog(QDialog):
    def __init__(self, versao, texto, zoom_atual=100, parent=None):
        super().__init__(parent)
        import idiomas
        self.setWindowTitle(idiomas.tr("nov_janela").format(versao))
        
        # MÁGICA: A largura se adapta ao zoom para não esmagar o texto gigante!
        nova_largura = int(650 * (zoom_atual / 100.0))
        self.setFixedSize(nova_largura, 650) # Altura cravada em 650, largura dinâmica
        
        layout = QVBoxLayout(self)
        
        lbl_titulo = QLabel(idiomas.tr("nov_header").format(versao))
        lbl_titulo.setStyleSheet("font-size: 16px;")
        layout.addWidget(lbl_titulo)
        
        # Cria a área de rolagem
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff) # <-- DESTRÓI A BARRA HORIZONTAL
        scroll.setStyleSheet("QScrollArea { border: 1px solid gray; border-radius: 5px; }")
        
        conteudo = QWidget()
        lay_conteudo = QVBoxLayout(conteudo)
        
        lbl_texto = QLabel(texto)
        lbl_texto.setWordWrap(True)
        lbl_texto.setAlignment(Qt.AlignTop | Qt.AlignLeft)
        
        lay_conteudo.addWidget(lbl_texto)
        lay_conteudo.addStretch() # Empurra o texto para o topo
        
        scroll.setWidget(conteudo)
        layout.addWidget(scroll)
        
        btn_ok = QPushButton(idiomas.tr("nov_btn_ok"))
        btn_ok.setStyleSheet("background-color: #0078d7; color: white; font-weight: bold; padding: 10px;")
        btn_ok.setCursor(Qt.PointingHandCursor)
        btn_ok.clicked.connect(self.accept)
        layout.addWidget(btn_ok)

# --- O MOTOR DO TUTORIAL (OVERLAY DE DESFOQUE) ---
class TutorialOverlay(QWidget):
    def __init__(self, parent):
        super().__init__(parent)
        self.passos = []
        self.passo_atual = 0
        
        # CRUCIAL: Bloqueia cliques para a interface, e permite transparência verdadeira
        self.setAttribute(Qt.WA_TransparentForMouseEvents, False)
        self.setAttribute(Qt.WA_TranslucentBackground, True)
        self.hide()

        self.lbl_titulo = QLabel(self)
        self.lbl_titulo.setStyleSheet("color: #ffc107; font-weight: 900; font-size: 24px; background-color: transparent;")
        
        self.lbl_texto = QLabel(self)
        self.lbl_texto.setStyleSheet("color: white; font-size: 16px; font-weight: 500; background-color: transparent;")
        self.lbl_texto.setWordWrap(True)

        self.btn_next = QPushButton(self)
        self.btn_next.setCursor(Qt.PointingHandCursor)
        self.btn_next.clicked.connect(self.avancar)

        self.btn_close = QPushButton(self)
        self.btn_close.setStyleSheet("background-color: #dc3545; color: white; font-weight: bold; padding: 10px 20px; border-radius: 6px; font-size: 14px; border: 2px solid white;")
        self.btn_close.setCursor(Qt.PointingHandCursor)
        self.btn_close.clicked.connect(self.encerrar)

    def iniciar(self, passos):
        import idiomas
        self.btn_close.setText(idiomas.tr("tut_btn_sair"))
        self.passos = passos
        self.passo_atual = 0
        self.setGeometry(0, 0, self.parentWidget().width(), self.parentWidget().height())
        self.show()
        self.raise_() # Joga por cima de TUDO
        self.atualizar_passo()

    def avancar(self):
        self.passo_atual += 1
        if self.passo_atual >= len(self.passos):
            self.encerrar()
        else:
            self.atualizar_passo()

    def encerrar(self):
        self.hide()

    def atualizar_passo(self):
        alvo, titulo, texto = self.passos[self.passo_atual]
        
        import idiomas
        self.lbl_titulo.setText(titulo)
        self.lbl_texto.setText(texto)
        self.btn_next.setText(idiomas.tr("tut_btn_fim") if self.passo_atual == len(self.passos) - 1 else idiomas.tr("tut_btn_prox"))
        self.btn_next.setStyleSheet("background-color: #28a745; color: white; font-weight: bold; padding: 10px 20px; border-radius: 6px; font-size: 14px; border: 2px solid white;" if self.passo_atual == len(self.passos) - 1 else "background-color: #0078d7; color: white; font-weight: bold; padding: 10px 20px; border-radius: 6px; font-size: 14px; border: 2px solid white;")

        if alvo:
            # A MÁGICA 1: Rola a tela até o alvo sozinho!
            scroll = self.parentWidget().scroll_principal
            scroll.ensureWidgetVisible(alvo, 50, 50)
            QApplication.processEvents() # Dá tempo pro PySide6 calcular a barra de rolagem

        self.posicionar_ui()
        self.update() # Força o paintEvent (recorte do buraco)

    def posicionar_ui(self):
        self.lbl_texto.setFixedWidth(400)
        self.lbl_titulo.adjustSize()
        self.lbl_texto.adjustSize()
        self.btn_next.adjustSize()
        self.btn_close.adjustSize()

        alvo = self.passos[self.passo_atual][0]
        
        y_caixa = 80
        if alvo:
            pos_global = alvo.mapToGlobal(QPoint(0, 0))
            pos_alvo = self.mapFromGlobal(pos_global)
            # Se o alvo tá do meio pra cima, jogamos a explicação pra baixo. Vice-versa.
            if pos_alvo.y() < self.height() / 2:
                y_caixa = self.height() - self.lbl_texto.height() - 120
            else:
                y_caixa = 80

        cx = (self.width() - 400) // 2
        self.lbl_titulo.move(cx, y_caixa)
        self.lbl_texto.move(cx, y_caixa + self.lbl_titulo.height() + 15)
        
        y_botoes = self.lbl_texto.y() + self.lbl_texto.height() + 25
        self.btn_close.move(cx, y_botoes)
        self.btn_next.move(cx + self.btn_close.width() + 15, y_botoes)

    def paintEvent(self, event):
        from PySide6.QtGui import QPainterPath
        from PySide6.QtCore import Qt, QPoint, QRect
        
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)

        alvo = self.passos[self.passo_atual][0]
        rect = None

        if alvo:
            # Cálculo ABSOLUTO para ignorar o bug do scroll
            pos_global = alvo.mapToGlobal(QPoint(0, 0))
            pos = self.mapFromGlobal(pos_global)
            rect = QRect(pos, alvo.size())
            rect.adjust(-10, -10, 10, 10) # Margem de respiro pro widget brilhar

        # A MÁGICA DEFINITIVA: QPainterPath com regra OddEven
        path = QPainterPath()
        path.addRect(self.rect()) # Adiciona a tela toda
        if rect:
            path.addRoundedRect(rect, 8, 8) # Adiciona o buraco dentro

        # Preenche tudo com preto 85%, exceto o buraco!
        painter.fillPath(path, QColor(0, 0, 0, 220))

        # Desenha a borda neon por cima
        if rect:
            pen = QPen(QColor("#ffc107"), 3, Qt.DashLine)
            painter.setPen(pen)
            painter.drawRoundedRect(rect, 8, 8)

# --- APLICATIVO PRINCIPAL ---
class AppPaiVega(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("SmartSlides Pro")
        self.resize(650, 650)
        
        # --- ÍCONE DA JANELA E BARRA DE TAREFAS ---
        icon_path = os.path.join(bundle_dir, "assets", "SmartSlides.ico") # <-- CORRIGIDO PARA BUNDLE_DIR
        if os.path.exists(icon_path):
            self.setWindowIcon(QIcon(icon_path))
            # Macete do CTYPES: Força o Windows a respeitar o seu ícone
            if os.name == 'nt':
                myappid = 'vega.smartslides.pro.1.0'
                ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID(myappid)
        
        self.config = self.load_config()
        self.last_dir = self.config.get("last_dir", "") # GUARDA A ÚLTIMA PASTA
        self.pptx_path = None
        self.is_manual_pptx = False # <-- TRAVA ANTI-CONFLITO
        self.last_saved_pptx = None # Guarda o último editado para o conversor de PDF
        self.bg_path = None
        self.text_color_title = None # <-- COR SEPARADA DO TÍTULO
        self.text_color_body = None  # <-- COR SEPARADA DO CORPO
        self.fill_color = None

        # A MÁGICA ACONTECE AQUI: Aplica o zoom e o tema ANTES de desenhar a interface!
        self.apply_zoom(self.config.get("zoom", 100))
        self.apply_theme()
        
        import idiomas
        idiomas.set_idioma(self.config.get("idioma", "pt"))
        
        self.init_ui() # Agora todos os textos e botões já nascem do tamanho salvo!
        self.atualizar_textos_ui() # Aciona a tradução pela primeira vez
        
        # Instancia o buraco negro invisível por cima da UI
        self.overlay_tutorial = TutorialOverlay(self)
        
        # Verifica se o atualizador deixou algum recado de novidades
        self.checar_notas_atualizacao()

    def load_config(self):
        import base64
        import json
        default_config = {
            "license": "FREE", 
            "api_key": "", 
            "model": "gemini-2.5-flash", 
            "theme": "Escuro", 
            "zoom": 100,
            "idioma": "pt",
            "saas_elite_ativo": False,
            "saas_padrao_ativo": False,
            "byok_ativo": False
        }
        
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    conteudo = f.read().strip()
                    
                if conteudo.startswith("{"): 
                    dados = json.loads(conteudo)
                    if dados.get("license") == "SAAS (Chave Embutida)" and not dados.get("token_master"):
                        dados["license"] = "FREE" 
                    # Garante que os campos de ativação venham preenchidos se existirem
                    for k, v in default_config.items():
                        if k not in dados: dados[k] = v
                    return dados
                else:
                    texto_reverso = conteudo[::-1]
                    json_str = base64.b64decode(texto_reverso.encode('utf-8')).decode('utf-8')
                    dados = json.loads(json_str)
                    for k, v in default_config.items():
                        if k not in dados: dados[k] = v
                    return dados
            except: pass
            
        return default_config

    def save_config(self):
        import base64
        import json
        self.config["last_dir"] = self.last_dir
        try:
            json_str = json.dumps(self.config)
            
            # GOLPE VEGATECH: Aplica Base64 e inverte a string pra confundir xeretas
            b64_str = base64.b64encode(json_str.encode('utf-8')).decode('utf-8')
            texto_ofuscado = b64_str[::-1]
            
            with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
                f.write(texto_ofuscado)
        except Exception as e:
            print(f"Erro na blindagem do config: {e}")

    def open_config(self):
        dialog = ConfigDialog(self.config, self)
        if dialog.exec() == QDialog.Accepted:
            self.config = dialog.config
            self.save_config()
            self.apply_theme()
            self.apply_zoom(self.config["zoom"])
            self.update_main_ui_lock()

    def apply_theme(self):
        tema = self.config.get("theme", "Escuro")
        app = QApplication.instance()
        app.setStyle("Fusion")
        
        palette = QPalette()
        if tema == "Escuro":
            # Cores Clássicas do Tema Escuro
            palette.setColor(QPalette.Window, QColor(53, 53, 53))
            palette.setColor(QPalette.WindowText, Qt.white)
            palette.setColor(QPalette.Base, QColor(25, 25, 25))
            palette.setColor(QPalette.AlternateBase, QColor(53, 53, 53))
            palette.setColor(QPalette.ToolTipBase, Qt.black)
            palette.setColor(QPalette.ToolTipText, Qt.white)
            palette.setColor(QPalette.Text, Qt.white)
            palette.setColor(QPalette.Button, QColor(53, 53, 53))
            palette.setColor(QPalette.ButtonText, Qt.white)
            palette.setColor(QPalette.BrightText, Qt.red)
            palette.setColor(QPalette.Link, QColor(42, 130, 218))
            palette.setColor(QPalette.Highlight, QColor(42, 130, 218))
            palette.setColor(QPalette.HighlightedText, Qt.black)
        else:
            # Tema CLARO - Contraste Absoluto e Legível
            palette.setColor(QPalette.Window, QColor(240, 240, 240))
            palette.setColor(QPalette.WindowText, Qt.black)
            palette.setColor(QPalette.Base, Qt.white)
            palette.setColor(QPalette.AlternateBase, QColor(225, 225, 225))
            palette.setColor(QPalette.ToolTipBase, Qt.white)
            palette.setColor(QPalette.ToolTipText, Qt.black)
            palette.setColor(QPalette.Text, Qt.black)
            palette.setColor(QPalette.Button, QColor(240, 240, 240))
            palette.setColor(QPalette.ButtonText, Qt.black)
            palette.setColor(QPalette.BrightText, Qt.red)
            palette.setColor(QPalette.Link, QColor(0, 102, 204))
            palette.setColor(QPalette.Highlight, QColor(0, 120, 215))
            palette.setColor(QPalette.HighlightedText, Qt.white)
            
        app.setPalette(palette)
        
        # Injeta um estilo blindado para que os Checkboxes fiquem sempre desenhados e visíveis
        estilo_checkbox = """
            QCheckBox::indicator {
                width: 16px; height: 16px; border: 2px solid #777; border-radius: 3px; background-color: transparent;
            }
            QCheckBox::indicator:checked {
                background-color: #0078d7; border-color: #0078d7;
            }
        """
        app.setStyleSheet(estilo_checkbox)

        # Aplica a paleta apenas no TopLevel para não estragar a herança interna dos botões
        for w in app.topLevelWidgets():
            w.setPalette(palette)
            w.update()

    def apply_zoom(self, val):
        # 1. Fonte: Começa em 11 e aumenta conforme o slider
        base_size = 11
        # O slider vai de 0 a 100 (ou mais). 0 = 100% do tamanho, 50 = 150% do tamanho
        fator = 1.0 + (val / 100.0) 
        new_size = int(base_size * fator)
        
        app = QApplication.instance()
        font = app.font()
        font.setPointSize(new_size)
        app.setFont(font)
        self.setFont(font)
        
        # 2. Janela: Base de 650px, aumenta conforme o slider
        nova_largura = int(650 * fator)
        self.setMinimumWidth(nova_largura)
        self.resize(nova_largura, self.height())
        
        # Garante que qualquer outra janela filha atualize na mesma hora
        for w in app.topLevelWidgets():
            w.setFont(font)
            w.update()

    def create_color_indicator(self):
        lbl = QLabel()
        lbl.setFixedSize(24, 24)
        lbl.setStyleSheet("background-color: transparent; border: 1px solid gray; border-radius: 12px;")
        return lbl

    def init_ui(self):
        self.main_container = QWidget()
        self.main_layout = QVBoxLayout(self.main_container)
        self.main_layout.setSpacing(15)

        hbox_top = QHBoxLayout()
        hbox_top.addStretch()
        
        self.btn_tutorial = QPushButton("📚 Como Usar")
        self.btn_tutorial.setStyleSheet("background-color: #0078d7; color: white; font-weight: bold; padding: 5px;")
        self.btn_tutorial.setCursor(Qt.PointingHandCursor)
        self.btn_tutorial.clicked.connect(self.abrir_tutorial)
        hbox_top.addWidget(self.btn_tutorial)
        
        self.btn_config = QPushButton("⚙️ Configurações / Licença")
        self.btn_config.setCursor(Qt.PointingHandCursor)
        self.btn_config.clicked.connect(self.open_config)
        hbox_top.addWidget(self.btn_config)
        
        # --- NOVO: BOTÃO PARA LER O RELEASE NOTES A QUALQUER MOMENTO ---
        self.btn_novidades = QPushButton("📋 Novidades")
        self.btn_novidades.setStyleSheet("background-color: #28a745; color: white; font-weight: bold; padding: 5px;")
        self.btn_novidades.setCursor(Qt.PointingHandCursor)
        self.btn_novidades.clicked.connect(self.abrir_novidades)
        hbox_top.addWidget(self.btn_novidades)
        # ---------------------------------------------------------------
        
        self.main_layout.addLayout(hbox_top)

        self.btn_load = QPushButton("1. Carregar Apresentação Pronta (.pptx)")
        self.btn_load.clicked.connect(self.load_pptx)
        self.lbl_pptx = QLabel("Nenhum arquivo selecionado")
        self.lbl_pptx.setStyleSheet("color: #d9534f; font-weight: bold; font-size: 15px;") # Vermelho, gordo e grande

        self.main_layout.addWidget(self.btn_load)
        self.main_layout.addWidget(self.lbl_pptx)

        from PySide6.QtWidgets import QSizePolicy
        from PySide6.QtCore import QEvent

        hbox_ia = QHBoxLayout()
        
        # --- BLOCO DA IA COM TÍTULO E GATILHO ---
        vbox_tema = QVBoxLayout()
        self.lbl_titulo_ia = QLabel()
        self.lbl_titulo_ia.setStyleSheet("font-weight: bold; font-size: 14px; color: #0078d7;")
        
        self.txt_tema = QTextEdit()
        self.txt_tema.setPlaceholderText("Digite o tema para a IA criar a apresentação...")
        self.txt_tema.setMaximumHeight(80) 
        # O SEGREDO DO QTEXTEDIT: Tem que escutar o VIEWPORT, não o widget solto!
        self.txt_tema.viewport().installEventFilter(self) 
        
        vbox_tema.addWidget(self.lbl_titulo_ia)
        vbox_tema.addWidget(self.txt_tema)
        
        self.btn_gerar_ia = QPushButton("Criar Apresentação\nAutomática (IA)")
        self.btn_gerar_ia.setStyleSheet("background-color: #2b5c8f; color: white; font-weight: bold; padding: 10px;")
        self.btn_gerar_ia.setSizePolicy(QSizePolicy.Minimum, QSizePolicy.MinimumExpanding) # Libera o crescimento
        self.btn_gerar_ia.clicked.connect(self.gerar_com_ia)
        
        self.btn_historico = QPushButton("🕒 Recuperar\nSalvo")
        self.btn_historico.setStyleSheet("background-color: #555; color: white; font-weight: bold; padding: 10px;")
        self.btn_historico.setSizePolicy(QSizePolicy.Minimum, QSizePolicy.MinimumExpanding) # Libera o crescimento
        self.btn_historico.setCursor(Qt.PointingHandCursor)
        self.btn_historico.clicked.connect(self.abrir_historico)
        
        hbox_ia.addLayout(vbox_tema)
        hbox_ia.addWidget(self.btn_gerar_ia)
        hbox_ia.addWidget(self.btn_historico)
        self.main_layout.addLayout(hbox_ia)
        
        # --- DROPDOWN DE PRESETS (ITEM 5) ---
        hbox_preset = QHBoxLayout()
        self.combo_presets = QComboBox()
        self.combo_presets.addItem("Personalizado (Manual)")
        self.combo_presets.addItem("Corporativo Azul")
        self.combo_presets.addItem("Dark Mode Clássico")
        self.combo_presets.addItem("Minimalista Clean")
        self.combo_presets.currentIndexChanged.connect(self.aplicar_preset)
        self.combo_presets.setCursor(Qt.PointingHandCursor)
        self.combo_presets.setStyleSheet("padding: 4px;") # Arrancado o font-size fixo!
        
        self.lbl_presets = QLabel()
        hbox_preset.addWidget(self.lbl_presets)
        hbox_preset.addWidget(self.combo_presets)
        hbox_preset.addStretch()
        self.main_layout.addLayout(hbox_preset)

        # --- 2. SELEÇÃO DE COR DE TEXTO SEPARADA (ITEM 3) ---
        hbox_color_title = QHBoxLayout()
        self.btn_color_title = QPushButton("🎨 Cor dos Títulos")
        self.btn_color_title.setCursor(Qt.PointingHandCursor)
        self.btn_color_title.clicked.connect(self.choose_title_color)
        self.lbl_color_title_indicator = self.create_color_indicator()
        
        self.lbl_cor_titulos = QLabel()
        hbox_color_title.addWidget(self.lbl_cor_titulos)
        hbox_color_title.addWidget(self.btn_color_title)
        hbox_color_title.addWidget(self.lbl_color_title_indicator)
        hbox_color_title.addStretch()
        self.main_layout.addLayout(hbox_color_title)

        hbox_color_body = QHBoxLayout()
        self.btn_color_body = QPushButton()
        self.btn_color_body.setCursor(Qt.PointingHandCursor)
        self.btn_color_body.clicked.connect(self.choose_body_color)
        self.lbl_color_body_indicator = self.create_color_indicator()
        
        self.lbl_cor_textos = QLabel()
        hbox_color_body.addWidget(self.lbl_cor_textos)
        hbox_color_body.addWidget(self.btn_color_body)
        hbox_color_body.addWidget(self.lbl_color_body_indicator)
        hbox_color_body.addStretch()
        self.main_layout.addLayout(hbox_color_body)
        
        # --- SELEÇÃO DE FONTE COM PREVIEW (SUGESTÃO 4) ---
        # --- SELEÇÃO DE FONTES (ITEM 2) ---
        fontes_comuns = ["Padrão do Tema", "Arial", "Calibri", "Century Gothic", "Georgia", "Segoe UI", "Tahoma", "Times New Roman", "Trebuchet MS", "Verdana"]
        
        # FONTE DOS TÍTULOS
        hbox_font_title = QHBoxLayout()
        self.combo_font_title = QComboBox()
        for i, nome_fonte in enumerate(fontes_comuns):
            self.combo_font_title.addItem(nome_fonte, nome_fonte)
            if i > 0: # 0 é o Padrão
                self.combo_font_title.setItemData(i, QFont(nome_fonte), Qt.FontRole)
        self.combo_font_title.setCursor(Qt.PointingHandCursor)

        self.lbl_font_title_preview = QLabel("Aa Título")
        self.lbl_font_title_preview.setStyleSheet("font-size: 18px; color: gray; font-weight: bold;")
        self.combo_font_title.currentIndexChanged.connect(
            lambda idx: self.lbl_font_title_preview.setStyleSheet(f"font-size: 18px; font-family: '{self.combo_font_title.itemData(idx)}'; font-weight: bold; color: #0078d7;") if idx != 0 else self.lbl_font_title_preview.setStyleSheet("font-size: 18px; color: gray; font-weight: bold;")
        )
        
        self.btn_add_font = QPushButton("➕ Nova Fonte")
        self.btn_add_font.setCursor(Qt.PointingHandCursor)
        self.btn_add_font.clicked.connect(self.adicionar_fonte_customizada)
        self.btn_add_font.setStyleSheet("padding: 4px; font-weight: bold;")

        self.lbl_fonte_titulos = QLabel()
        hbox_font_title.addWidget(self.lbl_fonte_titulos)
        hbox_font_title.addWidget(self.combo_font_title)
        hbox_font_title.addWidget(self.btn_add_font)
        hbox_font_title.addWidget(self.lbl_font_title_preview)
        hbox_font_title.addStretch()
        self.main_layout.addLayout(hbox_font_title)

        # FONTE DOS TEXTOS (CORPO)
        hbox_font_body = QHBoxLayout()
        self.combo_font_body = QComboBox()
        for i, nome_fonte in enumerate(fontes_comuns):
            self.combo_font_body.addItem(nome_fonte, nome_fonte)
            if i > 0:
                self.combo_font_body.setItemData(i, QFont(nome_fonte), Qt.FontRole)
        self.combo_font_body.setCursor(Qt.PointingHandCursor)
                
        # Avisa o mini-slide quando o usuário trocar as fontes
        self.combo_font_title.currentIndexChanged.connect(self.atualizar_preview_real)
        self.combo_font_body.currentIndexChanged.connect(self.atualizar_preview_real)

        self.lbl_font_body_preview = QLabel("Aa Corpo")
        self.lbl_font_body_preview.setStyleSheet("font-size: 14px; color: gray;")
        self.combo_font_body.currentIndexChanged.connect(
            lambda idx: self.lbl_font_body_preview.setStyleSheet(f"font-size: 14px; font-family: '{self.combo_font_body.itemData(idx)}'; color: #0078d7;") if idx != 0 else self.lbl_font_body_preview.setStyleSheet("font-size: 14px; color: gray;")
        )
        
        self.lbl_fonte_textos = QLabel()
        hbox_font_body.addWidget(self.lbl_fonte_textos)
        hbox_font_body.addWidget(self.combo_font_body)
        hbox_font_body.addWidget(self.lbl_font_body_preview)
        hbox_font_body.addStretch()
        self.main_layout.addLayout(hbox_font_body)

        hbox_fill = QHBoxLayout()
        self.btn_fill_color = QPushButton("3. Escolher Cor de Fundo dos Cards (Opcional)")
        self.btn_fill_color.clicked.connect(self.choose_fill_color)
        self.lbl_fill_color_indicator = self.create_color_indicator()
        
        hbox_fill.addWidget(self.btn_fill_color)
        hbox_fill.addWidget(self.lbl_fill_color_indicator)
        hbox_fill.addStretch()
        self.main_layout.addLayout(hbox_fill)

        hbox_bg = QHBoxLayout()
        self.btn_bg = QPushButton("4. Fundo do Slide (Opcional)")
        self.btn_bg.clicked.connect(self.load_bg)
        self.lbl_bg_preview = QLabel("Sem imagem")
        self.lbl_bg_preview.setFixedSize(80, 80)
        self.lbl_bg_preview.setAlignment(Qt.AlignCenter)
        self.lbl_bg_preview.setStyleSheet("border: 1px dashed gray;")
        
        hbox_bg.addWidget(self.btn_bg)
        hbox_bg.addWidget(self.lbl_bg_preview)
        hbox_bg.addStretch()
        self.main_layout.addLayout(hbox_bg)

        # --- 5. ESCALA DE IMAGENS COM SLIDER (DIMINUI E AUMENTA - ITEM 2) ---
        hbox_scale = QHBoxLayout()
        self.chk_scale = QCheckBox("5. Escala das imagens:")
        self.slider_scale = QSlider(Qt.Horizontal)
        self.slider_scale.setRange(10, 200) # 10% até 200% do tamanho
        self.slider_scale.setValue(100) # 100% é o tamanho original exacto
        self.slider_scale.setEnabled(False)
        self.lbl_scale_val = QLabel("100%")
        self.lbl_scale_val.setStyleSheet("font-weight: bold;")
        
        self.chk_scale.toggled.connect(self.slider_scale.setEnabled)
        self.slider_scale.valueChanged.connect(lambda v: self.lbl_scale_val.setText(f"{v}%"))
        
        hbox_scale.addWidget(self.chk_scale)
        hbox_scale.addWidget(self.slider_scale)
        hbox_scale.addWidget(self.lbl_scale_val)
        self.main_layout.addLayout(hbox_scale)

        # --- 5.1. AUMENTO DE TEXTOS COM SLIDER ---
        hbox_text_scale = QHBoxLayout()
        self.chk_text_scale = QCheckBox("5.1. Aumentar textos em:")
        self.slider_text_scale = QSlider(Qt.Horizontal)
        self.slider_text_scale.setRange(1, 150)
        self.slider_text_scale.setValue(20)
        self.slider_text_scale.setEnabled(False)
        self.lbl_text_scale_val = QLabel("20%")
        self.lbl_text_scale_val.setStyleSheet("font-weight: bold;")
        
        self.chk_text_scale.toggled.connect(self.slider_text_scale.setEnabled)
        self.slider_text_scale.valueChanged.connect(lambda v: self.lbl_text_scale_val.setText(f"{v}%"))
        
        hbox_text_scale.addWidget(self.chk_text_scale)
        hbox_text_scale.addWidget(self.slider_text_scale)
        hbox_text_scale.addWidget(self.lbl_text_scale_val)
        self.main_layout.addLayout(hbox_text_scale)

        # --- PREVISÃO AO VIVO (MINI-SLIDE) ---
        hbox_preview = QHBoxLayout()
        self.frame_preview = QFrame()
        self.frame_preview.setFixedSize(320, 180) # Formato 16:9 perfeito
        self.frame_preview.setStyleSheet("background-color: white; border: 2px solid #ccc; border-radius: 8px;")
        
        # Elementos posicionados de forma absoluta no slide virtual
        self.lbl_preview_text = QLabel("Título Exemplo\n\nEste é o corpo do texto.\nVeja como ele cresce.", self.frame_preview)
        self.lbl_preview_text.setGeometry(10, 10, 180, 160)
        self.lbl_preview_text.setAlignment(Qt.AlignTop | Qt.AlignLeft)
        self.lbl_preview_text.setStyleSheet("color: #333; border: none; font-size: 10px;")

        self.lbl_preview_img = QLabel(self.frame_preview)
        self.lbl_preview_img.setStyleSheet("background-color: #2b5c8f; border-radius: 5px; border: none;")
        self.lbl_preview_img.setGeometry(190, 40, 100, 100) # Posição e tamanho base

        hbox_preview.addStretch()
        hbox_preview.addWidget(self.frame_preview)
        hbox_preview.addStretch()
        self.main_layout.addLayout(hbox_preview)
        
        # Liga os gatilhos dos sliders e checkboxes ao método da classe
        self.chk_scale.toggled.connect(self.atualizar_preview_real)
        self.slider_scale.valueChanged.connect(self.atualizar_preview_real)
        self.chk_text_scale.toggled.connect(self.atualizar_preview_real)
        self.slider_text_scale.valueChanged.connect(self.atualizar_preview_real)
        
        self.atualizar_preview_real() # Força o primeiro desenho

        # Liga os gatilhos dos sliders e checkboxes à função de animação
        self.chk_scale.toggled.connect(self.atualizar_preview_real)
        self.slider_scale.valueChanged.connect(self.atualizar_preview_real)
        self.chk_text_scale.toggled.connect(self.atualizar_preview_real)
        self.slider_text_scale.valueChanged.connect(self.atualizar_preview_real)
        
        self.atualizar_preview_real() # Força o primeiro desenho # Força o primeiro desenho
        
        # --- CHAVE PARA LIGAR/DESLIGAR AS ANOTAÇÕES (SUGESTÃO 2) ---
        hbox_notes = QHBoxLayout()
        self.chk_notes = QCheckBox("5.2. Inserir Roteiro do Apresentador nas anotações nativas")
        self.chk_notes.setChecked(True) # Deixamos ativo por padrão, o usuário desmarca se quiser
        
        self.chk_barra_lateral = QCheckBox("5.3. Exibir barra azul lateral")
        self.chk_barra_lateral.setChecked(True)
        
        hbox_notes.addWidget(self.chk_notes)
        hbox_notes.addWidget(self.chk_barra_lateral)
        hbox_notes.addStretch()
        self.main_layout.addLayout(hbox_notes)

        self.btn_save = QPushButton("6. Processar Apresentação e Salvar")
        self.btn_save.setSizePolicy(QSizePolicy.Minimum, QSizePolicy.MinimumExpanding)
        self.btn_save.setStyleSheet("background-color: #333; color: white; font-weight: bold; padding: 10px;")
        self.btn_save.clicked.connect(self.process_and_save)
        self.main_layout.addWidget(self.btn_save)

        # --- 7. BOTAO CONVERTER PARA PDF ---
        self.btn_pdf = QPushButton("7. Converter o PPTX atual para PDF")
        self.btn_pdf.setSizePolicy(QSizePolicy.Minimum, QSizePolicy.MinimumExpanding)
        self.btn_pdf.setStyleSheet("background-color: #a8201a; color: white; font-weight: bold; padding: 10px;")
        self.btn_pdf.clicked.connect(self.export_to_pdf)
        self.main_layout.addWidget(self.btn_pdf)

        # =========================================================
        # MÁGICA RESPONSIVA: Barra de Rolagem Inteligente
        # =========================================================
        from PySide6.QtWidgets import QScrollArea
        
        self.scroll_principal = QScrollArea()
        self.scroll_principal.setWidgetResizable(True) 
        self.scroll_principal.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff) # <-- ASSASSINA A BARRA HORIZONTAL
        self.scroll_principal.setStyleSheet("QScrollArea { border: none; background-color: transparent; }")
        self.scroll_principal.setWidget(self.main_container)
        
        # Agora o widget central da tela é a área de rolagem, que guarda todo o resto dentro dela
        self.setCentralWidget(self.scroll_principal)
        
        self.update_main_ui_lock()   

    def atualizar_textos_ui(self):
        import idiomas
        self.btn_tutorial.setText(idiomas.tr("btn_tutorial"))
        self.btn_config.setText(idiomas.tr("btn_config"))
        self.btn_novidades.setText(idiomas.tr("btn_novidades"))
        self.btn_load.setText(idiomas.tr("btn_load"))
        if not self.pptx_path: self.lbl_pptx.setText(idiomas.tr("lbl_pptx_vazio"))
        self.lbl_titulo_ia.setText(idiomas.tr("lbl_titulo_ia"))
        self.btn_gerar_ia.setText(idiomas.tr("btn_ia"))
        self.btn_historico.setText(idiomas.tr("btn_historico"))
        self.lbl_presets.setText(idiomas.tr("lbl_presets"))
        self.lbl_cor_titulos.setText(idiomas.tr("lbl_cor_titulos"))
        self.btn_color_title.setText(idiomas.tr("btn_cor_titulos_btn"))
        self.lbl_cor_textos.setText(idiomas.tr("lbl_cor_textos"))
        self.btn_color_body.setText(idiomas.tr("btn_cor_textos_btn"))
        self.btn_add_font.setText(idiomas.tr("btn_add_font"))
        self.lbl_fonte_titulos.setText(idiomas.tr("lbl_fonte_titulos"))
        self.lbl_fonte_textos.setText(idiomas.tr("lbl_fonte_textos"))
        self.btn_fill_color.setText(idiomas.tr("btn_fill_color"))
        self.btn_bg.setText(idiomas.tr("btn_bg"))
        if not self.bg_path: self.lbl_bg_preview.setText(idiomas.tr("lbl_bg_preview"))
        self.chk_scale.setText(idiomas.tr("chk_scale"))
        self.chk_text_scale.setText(idiomas.tr("chk_text_scale"))
        self.chk_notes.setText(idiomas.tr("chk_notes"))
        self.chk_barra_lateral.setText(idiomas.tr("chk_barra_lateral"))
        self.btn_save.setText(idiomas.tr("btn_save"))
        self.btn_pdf.setText(idiomas.tr("btn_pdf"))
        
        # Traduz as opções do ComboBox de Presets
        self.combo_presets.setItemText(0, idiomas.tr("preset_manual"))
        self.combo_presets.setItemText(1, idiomas.tr("preset_corp"))
        self.combo_presets.setItemText(2, idiomas.tr("preset_dark"))
        self.combo_presets.setItemText(3, idiomas.tr("preset_clean"))

        # Traduz a opção "Padrão do Tema" nos dois ComboBoxes de Fonte
        self.combo_font_title.setItemText(0, idiomas.tr("font_default"))
        self.combo_font_body.setItemText(0, idiomas.tr("font_default"))

        self.update_main_ui_lock()
        self.atualizar_preview_real()

    def update_main_ui_lock(self):
        import idiomas
        if "FREE" in self.config.get("license", "FREE"):
            self.txt_tema.setReadOnly(True)
            self.txt_tema.setPlaceholderText(idiomas.tr("txt_tema_lock"))
            self.txt_tema.setStyleSheet("background-color: #444; color: #888; border: 1px solid #333;")
            
            self.btn_gerar_ia.setEnabled(False)
            self.btn_gerar_ia.setStyleSheet("background-color: #555; color: #888; font-weight: bold; padding: 10px;")
        else:
            self.txt_tema.setReadOnly(False)
            self.txt_tema.setPlaceholderText(idiomas.tr("txt_tema_placeholder"))
            self.txt_tema.setStyleSheet("") 
            
            self.btn_gerar_ia.setEnabled(True)
            self.btn_gerar_ia.setStyleSheet("background-color: #2b5c8f; color: white; font-weight: bold; padding: 10px;")

    # --- O DETETIVE DE CLIQUES (AVISO PARA CONVERSÃO DE CLIENTE) ---
    def eventFilter(self, obj, event):
        from PySide6.QtCore import QEvent
        # Agora ele verifica se o clique aconteceu na área interna (viewport) do campo de texto
        if obj == self.txt_tema.viewport() and event.type() == QEvent.MouseButtonPress:
            if "FREE" in self.config.get("license", "FREE"):
                QMessageBox.warning(
                    self, 
                    "Recurso Premium Bloqueado", 
                    "A geração de apresentações por Inteligência Artificial é um recurso VIP!\n\n"
                    "Clique no botão '⚙️ Configurações / Licença' no topo da tela, escolha um dos nossos planos na Loja e libere a IA para trabalhar por você."
                )
                return True # Consome o evento e não deixa o clique entrar na caixa de texto
        return super().eventFilter(obj, event)
            
    def atualizar_preview_real(self, *args):
        # --- 1. CAPTURAR AS FONTES E CORES ATUAIS ---
        fonte_titulo = self.combo_font_title.currentData()
        fonte_corpo = self.combo_font_body.currentData()
        
        fam_titulo = f"font-family: '{fonte_titulo}';" if self.combo_font_title.currentIndex() != 0 else ""
        fam_corpo = f"font-family: '{fonte_corpo}';" if self.combo_font_body.currentIndex() != 0 else ""
        
        # --- 1. CAPTURAR AS CORES SEPARADAS (ITEM 3) ---
        cor_titulo = f"color: {self.text_color_title.name()};" if self.text_color_title else "color: #1e3c5a;"
        cor_corpo = f"color: {self.text_color_body.name()};" if self.text_color_body else "color: #333;"
        
        # --- 2. ESCALA DO TEXTO ---
        base_px = 10
        if self.chk_text_scale.isChecked():
            new_px = int(base_px * (1 + (self.slider_text_scale.value() / 100.0)))
        else:
            new_px = base_px
            
        import idiomas
        # O truque sujo: Usar HTML no widget nativo para separar Título do Corpo perfeitamente
        html_texto = f"""
        <div style="{fam_titulo} font-size: {new_px + 4}px; font-weight: bold; {cor_titulo} margin-bottom: 5px;">
            {idiomas.tr("lbl_preview_text_title")}
        </div>
        <div style="{fam_corpo} font-size: {new_px}px; {cor_corpo}">
            {idiomas.tr("lbl_preview_text_body")}
        </div>
        """
        self.lbl_preview_text.setText(html_texto)
        self.lbl_preview_text.setStyleSheet("background-color: transparent; border: none;")

        # --- 3. FUNDO DO SLIDE (IMAGEM OU BRANCO) ---
        self.frame_preview.setObjectName("PreviewSlide") # Nomeia o frame para a imagem não bugar
        if self.bg_path:
            bg_limpo = self.bg_path.replace("\\", "/")
            # Usamos border-image com o ID (#PreviewSlide) para esticar a foto no 16:9 perfeito
            self.frame_preview.setStyleSheet(f"#PreviewSlide {{ border-image: url('{bg_limpo}'); border: 2px solid #ccc; border-radius: 8px; }}")
        else:
            self.frame_preview.setStyleSheet("#PreviewSlide { background-color: white; border: 2px solid #ccc; border-radius: 8px; }")

        # --- 4. COR E ESCALA DO CARD/IMAGEM (O QUADRADO) ---
        cor_fundo = self.fill_color.name() if self.fill_color else "#2b5c8f"
        self.lbl_preview_img.setStyleSheet(f"background-color: {cor_fundo}; border-radius: 5px; border: none;")
        
        if self.chk_scale.isChecked():
            base_w, base_h = 100, 100
            fator = self.slider_scale.value() / 100.0 # <--- 0.5 encolhe pela metade, 1.5 cresce 50%
            new_w, new_h = int(base_w * fator), int(base_h * fator)
            centro_x, centro_y = 240, 90 
            self.lbl_preview_img.setGeometry(centro_x - (new_w // 2), centro_y - (new_h // 2), new_w, new_h)
        else:
            self.lbl_preview_img.setGeometry(190, 40, 100, 100)

    def load_pptx(self):
        # --- TRAVA REVERSA (Impede o cara de jogar a IA no lixo sem querer) ---
        if self.pptx_path and not getattr(self, 'is_manual_pptx', True):
            resp = QMessageBox.critical(
                self, 
                "⚠️ ALERTA: Apresentação da IA já carregada!", 
                "VOCÊ JÁ GEROU UMA APRESENTAÇÃO COM A INTELIGÊNCIA ARTIFICIAL!\n\nSe você carregar um arquivo antigo do seu computador agora, VAI APAGAR a apresentação nova que a IA acabou de fazer.\n\nTem CERTEZA ABSOLUTA que quer jogar a da IA no lixo e carregar outro arquivo?", 
                QMessageBox.Yes | QMessageBox.No
            )
            if resp == QMessageBox.No:
                return # Aborta e salva a vida do usuário
        
        path, _ = QFileDialog.getOpenFileName(self, "Selecionar PPTX", self.last_dir, "PowerPoint (*.pptx)")
        if path:
            self.last_dir = os.path.dirname(path)
            self.save_config()
            self.pptx_path = path
            self.is_manual_pptx = True # <-- ACIONA O ALARME DE CONFLITO
            self.lbl_pptx.setText(os.path.basename(path))

    def choose_title_color(self):
        color = QColorDialog.getColor()
        if color.isValid():
            self.text_color_title = color
            self.lbl_color_title_indicator.setStyleSheet(f"background-color: {color.name()}; border: 1px solid gray; border-radius: 12px;")
            self.atualizar_preview_real()

    def choose_body_color(self):
        color = QColorDialog.getColor()
        if color.isValid():
            self.text_color_body = color
            self.lbl_color_body_indicator.setStyleSheet(f"background-color: {color.name()}; border: 1px solid gray; border-radius: 12px;")
            self.atualizar_preview_real()

    def choose_fill_color(self):
        color = QColorDialog.getColor()
        if color.isValid():
            self.fill_color = color
            self.lbl_fill_color_indicator.setStyleSheet(f"background-color: {color.name()}; border: 1px solid gray; border-radius: 12px;")
            self.atualizar_preview_real() # Avisa o mini-slide imediatamente
            
    def load_bg(self):
        path, _ = QFileDialog.getOpenFileName(self, "Selecionar Imagem de Fundo", self.last_dir, "Imagens (*.png *.jpg *.jpeg)")
        if path:
            self.last_dir = os.path.dirname(path)
            self.save_config()
            self.bg_path = path
            pixmap = QPixmap(path)
            self.lbl_bg_preview.setPixmap(pixmap.scaled(80, 80, Qt.KeepAspectRatio, Qt.SmoothTransformation))
            self.lbl_bg_preview.setStyleSheet("border: 1px solid gray;") 
            self.atualizar_preview_real() # Avisa o mini-slide imediatamente
            
    def aplicar_preset(self):
        idx = self.combo_presets.currentIndex()
        if idx == 0: return # Personalizado
            
        if idx == 1: # Corporativo Azul
            self.text_color_title = QColor(30, 60, 90)
            self.text_color_body = QColor(40, 50, 60)
            self.fill_color = QColor(240, 248, 255)
        elif idx == 2: # Dark Mode Clássico
            self.text_color_title = QColor(255, 255, 255)
            self.text_color_body = QColor(220, 220, 220)
            self.fill_color = QColor(45, 45, 45)
        elif idx == 3: # Minimalista Clean
            self.text_color_title = QColor(0, 0, 0)
            self.text_color_body = QColor(50, 50, 50)
            self.fill_color = QColor(255, 255, 255)
            
        self.lbl_color_title_indicator.setStyleSheet(f"background-color: {self.text_color_title.name()}; border: 1px solid gray; border-radius: 12px;")
        self.lbl_color_body_indicator.setStyleSheet(f"background-color: {self.text_color_body.name()}; border: 1px solid gray; border-radius: 12px;")
        self.lbl_fill_color_indicator.setStyleSheet(f"background-color: {self.fill_color.name()}; border: 1px solid gray; border-radius: 12px;")
        self.atualizar_preview_real()
            
    def adicionar_fonte_customizada(self):
        ok, font = QFontDialog.getFont(self)
        if ok:
            nome_fonte = font.family()
            # Adiciona a fonte escolhida nos DOIS dropdowns de uma vez
            for combo in [self.combo_font_title, self.combo_font_body]:
                index = combo.findText(nome_fonte)
                if index == -1:
                    combo.addItem(nome_fonte)
                    novo_index = combo.count() - 1
                    combo.setItemData(novo_index, QFont(nome_fonte), Qt.FontRole)
                    combo.setCurrentIndex(novo_index)
                else:
                    combo.setCurrentIndex(index)

    def gerar_com_ia(self):
        # --- PROTEÇÃO ANTI-CONFLITO DE ARQUIVOS ---
        if self.pptx_path and getattr(self, 'is_manual_pptx', False):
            resp = QMessageBox.warning(
                self, 
                "Atenção: Substituir Arquivo?", 
                "Você já carregou uma apresentação pronta no aplicativo (Botão 1).\n\nSe você usar a IA agora, a apresentação que você carregou será descartada e substituída por uma nova.\n\nDeseja descartar a apresentação atual e continuar com a IA?", 
                QMessageBox.Yes | QMessageBox.No
            )
            if resp == QMessageBox.No:
                return # Aborta e deixa o cara quieto com o arquivo dele
            else:
                # Limpa a trava e a tela para a IA trabalhar do zero
                self.pptx_path = None
                self.is_manual_pptx = False
                self.lbl_pptx.setText("Nenhum arquivo selecionado")
                
        licensa = self.config.get("license", "FREE")
        
        if licensa == "FREE":
            QMessageBox.warning(self, "Recurso Premium", "A geração por IA é um recurso Premium.")
            return

        tema = self.txt_tema.toPlainText().strip()
        if not tema:
            QMessageBox.warning(self, "Aviso", "Por favor, introduza o tema primeiro!")
            return
        
        api_key_usuario = self.config.get("api_key") if licensa == "BYOK (Sua Chave)" else None
        modelo = self.config.get("model", "gemini-2.5-flash")
        
        sucesso_check, msg_check = motor_ia.testar_conectividade(api_key_usuario)
        if not sucesso_check:
            QMessageBox.critical(self, "Bloqueio na IA", f"Conexão recusada:\n\n{msg_check}")
            return
            
        from PySide6.QtCore import QEventLoop

        try:
            progress = QProgressDialog("Inicializando Motor...", "Cancelar", 0, 100, self)
            progress.setWindowTitle("Gerando Apresentação")
            progress.setWindowModality(Qt.WindowModal)
            progress.setAutoClose(True)
            progress.show()
            
            # Trava os botões para ele não clicar duas vezes sem querer
            self.btn_gerar_ia.setEnabled(False)
            self.btn_historico.setEnabled(False)

            # Inicia o Trabalho Paralelo entregando todas as configs de IA
            worker = WorkerIAGerador(tema, self.config)
            # Liga o progresso da thread direto na barra da tela
            worker.progresso.connect(lambda v, t: (progress.setValue(v), progress.setLabelText(t)))
            progress.canceled.connect(worker.requestInterruption)
            
            loop = QEventLoop()
            worker.concluido.connect(loop.quit)
            worker.start()
            
            # A MÁGICA: A UI não congela mais, e o código espera o Worker terminar aqui!
            loop.exec() 
            
            self.btn_gerar_ia.setEnabled(True)
            self.btn_historico.setEnabled(True)
            progress.close()

            if worker.isInterruptionRequested(): return
                
            if worker.erro_msg:
                QMessageBox.critical(self, "Falha na IA", f"Erro:\n\n{worker.erro_msg}")
                return
                
            roteiro = worker.roteiro
            slide_images_dict = worker.slide_images_dict

            selected_images = {}
            if slide_images_dict:
                dialog = ImageSelectionDialog(slide_images_dict, roteiro, self)
                if dialog.exec() == QDialog.Accepted:
                    selected_images = dialog.selected_images
                else:
                    try:
                        for caminhos in slide_images_dict.values():
                            for c in caminhos:
                                if os.path.exists(c): os.remove(c)
                    except: pass
                    return

            prs = Presentation()
            prs.slide_width = 12192000
            prs.slide_height = 6858000

            for i, slide_data in enumerate(roteiro):
                slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                
                barra_lateral = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(prs.slide_width * 0.02), prs.slide_height)
                barra_lateral.fill.solid()
                barra_lateral.fill.fore_color.rgb = RGBColor(43, 92, 143) 
                barra_lateral.line.fill.background() 
                
                barra_lateral.name = "BarraLateral_Vega" # ETIQUETA PARA CONTROLE DE LIGA/DESLIGA

                # --- LÓGICA DINÂMICA DE LARGURA (ITEM 3) ---
                # --- LEITURA DO DESIGN DA IA ---
                tipo_layout = slide_data.get("tipo_layout", "padrao")
                # Se for cards ou destaque, a gente sabota a imagem para o design não quebrar
                tem_imagem = i in selected_images and selected_images[i] != "NONE" and tipo_layout == "padrao"
                largura_texto = int(prs.slide_width * 0.45) if tem_imagem else int(prs.slide_width * 0.88)

                # --- CAIXA DE TÍTULO (ITEM 4 - Esticado no slide) ---
                largura_titulo = int(prs.slide_width * 0.88) # Ignora a lateral e abraça 88% do topo
                top_titulo = int(prs.slide_height * 0.35) if tipo_layout == "destaque" else int(prs.slide_height * 0.06)
                altura_titulo = int(prs.slide_height * 0.30) if tipo_layout == "destaque" else int(prs.slide_height * 0.15)
                
                title_box = slide.shapes.add_textbox(int(prs.slide_width * 0.06), top_titulo, largura_titulo, altura_titulo)
                title_box.name = "TitleBox_Vega"
                tf_title = title_box.text_frame
                tf_title.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP # ALINHA NO TETO
                tf_title.word_wrap = True
                tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
                
                p_title = tf_title.paragraphs[0]
                p_title.text = slide_data.get("titulo", "").upper()
                p_title.font.bold = True
                p_title.font.size = Pt(50) if tipo_layout == "destaque" else Pt(28)
                p_title.font.color.rgb = RGBColor(30, 60, 90)

                # --- O MOTOR DE LAYOUTS ---
                if tipo_layout == "cards" and "topicos" in slide_data:
                    num_cards = min(len(slide_data["topicos"]), 3) # Trava em 3 para não virar bagunça
                    if num_cards > 0:
                        largura_card = int(largura_texto / num_cards) - int(prs.slide_width * 0.02)
                        altura_card = int(prs.slide_height * 0.55)
                        top_card = int(prs.slide_height * 0.24)
                        left_inicial = int(prs.slide_width * 0.06)

                        for j, topico in enumerate(slide_data["topicos"][:num_cards]):
                            card_left = left_inicial + (j * (largura_card + int(prs.slide_width * 0.02)))
                            # Desenha o Retângulo Arredondado do Card
                            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, top_card, largura_card, altura_card)
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
                            shape.line.color.rgb = RGBColor(200, 210, 220)

                            tf_card = shape.text_frame
                            tf_card.word_wrap = True
                            tf_card.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            
                            p_card = tf_card.paragraphs[0]
                            icone = topico.get("icone", "📌")
                            p_card.text = f"{icone} {topico.get('titulo', '')}"
                            p_card.font.bold = True
                            p_card.font.size = Pt(20)
                            p_card.font.color.rgb = RGBColor(43, 92, 143)

                            p_card_texto = tf_card.add_paragraph()
                            p_card_texto.text = topico.get("texto", "")
                            p_card_texto.font.size = Pt(14)
                            p_card_texto.font.color.rgb = RGBColor(60, 60, 60)
                            p_card_texto.space_before = Pt(10)
                            
                elif tipo_layout == "grafico" and "dados_grafico" in slide_data:
                    dados = slide_data.get("dados_grafico", {})
                    categorias = dados.get("categorias", [])
                    valores = dados.get("valores", [])
                    nome_serie = dados.get("nome_serie", "Dados")

                    if categorias and valores and len(categorias) == len(valores):
                        chart_data = CategoryChartData()
                        chart_data.categories = categorias
                        chart_data.add_series(nome_serie, valores)

                        x, y = int(prs.slide_width * 0.06), int(prs.slide_height * 0.24)
                        cx, cy = int(prs.slide_width * 0.88), int(prs.slide_height * 0.60)
                        
                        slide.shapes.add_chart(
                            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                        )

                elif tipo_layout != "destaque":
                    # Layout Padrão (Agora com suporte a Bullet Points e Ícones da IA)
                    body_box = slide.shapes.add_textbox(int(prs.slide_width * 0.06), int(prs.slide_height * 0.24), largura_texto, int(prs.slide_height * 0.60))
                    body_box.name = "BodyBox_Vega" # Etiqueta invisível do corpo
                    tf_body = body_box.text_frame
                    tf_body.word_wrap = True
                    tf_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 

                    texto_base = slide_data.get("texto", "").strip()
                    if texto_base:
                        p_body = tf_body.paragraphs[0]
                        p_body.text = texto_base
                        p_body.font.size = Pt(18)
                        p_body.space_after = Pt(14)

                    if "topicos" in slide_data:
                        for idx, topico in enumerate(slide_data["topicos"]):
                            p_topic = tf_body.add_paragraph() if texto_base else (tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph())
                            icone = topico.get("icone", "•")
                            p_topic.text = f"{icone} {topico.get('titulo', '')}: {topico.get('texto', '')}"
                            p_topic.font.size = Pt(16)
                            p_topic.level = 0
                            p_topic.space_after = Pt(10)

                # --- INSERÇÃO DA IMAGEM ---
                if tem_imagem:
                    try:
                        img_left = int(prs.slide_width * 0.54)
                        img_top = int(prs.slide_height * 0.24) # Desceu de 0.1 para 0.22 para fugir do título
                        img_max_w = int(prs.slide_width * 0.42)
                        img_max_h = int(prs.slide_height * 0.68)

                        pic = slide.shapes.add_picture(selected_images[i], img_left, img_top)
                        ratio = min(img_max_w / pic.width, img_max_h / pic.height)

                        pic.width = int(pic.width * ratio)
                        pic.height = int(pic.height * ratio)
                        pic.left = img_left + int((img_max_w - pic.width) / 2)
                        pic.top = img_top + int((img_max_h - pic.height) / 2)
                    except: pass
                    
                # --- ANOTAÇÕES DO APRESENTADOR (SUGESTÃO 2 com trava) ---
                if self.chk_notes.isChecked():
                    notas_palestrante = slide_data.get("roteiro_apresentador", "").strip()
                    if notas_palestrante:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        text_frame.text = f"📢 ROTEIRO DO APRESENTADOR:\n\n{notas_palestrante}"

            temp_path = TEMP_PPTX_FILE
            prs.save(temp_path)
            
            # --- SALVANDO O PPTX FÍSICO NO HISTÓRICO ---
            try:
                import shutil
                from datetime import datetime
                if not os.path.exists(PASTA_HISTORICO_PPTX):
                    os.makedirs(PASTA_HISTORICO_PPTX)
                
                # Cria um nome limpo e seguro pro arquivo
                tema_limpo = "".join([c for c in tema if c.isalpha() or c.isdigit() or c==' ']).rstrip()[:30]
                data_str = datetime.now().strftime("%Y%m%d_%H%M%S")
                nome_arq = f"{data_str}_{tema_limpo.replace(' ', '_')}.pptx"
                caminho_salvo = os.path.join(PASTA_HISTORICO_PPTX, nome_arq)
                
                shutil.copy2(temp_path, caminho_salvo)
                
                # Salva a referência no JSON blindado
                historico = []
                if os.path.exists(HISTORICO_JSON_FILE):
                    with open(HISTORICO_JSON_FILE, "r", encoding="utf-8") as f:
                        historico = json.load(f)
                        
                historico.append({
                    "data": datetime.now().strftime("%d/%m/%Y %H:%M:%S"),
                    "tema": tema,
                    "caminho": caminho_salvo
                })
                
                # MANTÉM APENAS OS 5 ÚLTIMOS E FAZ A FAXINA FÍSICA NO HD
                while len(historico) > 5:
                    removido = historico.pop(0)
                    try:
                        if os.path.exists(removido["caminho"]):
                            os.remove(removido["caminho"])
                    except: pass
                    
                with open(HISTORICO_JSON_FILE, "w", encoding="utf-8") as f:
                    json.dump(historico, f, indent=4, ensure_ascii=False)
            except Exception as e:
                print(f"Erro no backup do PPTX: {e}")
            
            try:
                for caminhos in slide_images_dict.values():
                    for c in caminhos:
                        if os.path.exists(c): os.remove(c)
            except: pass

            self.pptx_path = temp_path
            self.lbl_pptx.setText("✨ Apresentação carregada e pronta para edição final!")
            
            try:
                if os.name == 'nt': os.startfile(temp_path)
                else: subprocess.Popen(['open' if sys.platform == 'darwin' else 'xdg-open', temp_path])
            except: pass

        except Exception as e:
            progress.close()
            erro_msg = traceback.format_exc()
            
            with open(LOG_FILE, "w", encoding="utf-8") as f:
                f.write(f"ERRO CRÍTICO NA IA:\n{erro_msg}")
            
            QMessageBox.critical(self, "Falha Fatal", f"O aplicativo encontrou um erro e gerou um log.\n\nSalvo em:\n{LOG_FILE}\n\nErro: {str(e)}")
            
    def abrir_historico(self):
        if not os.path.exists(HISTORICO_JSON_FILE):
            QMessageBox.information(self, "Histórico Vazio", "Você ainda não tem apresentações salvas localmente.")
            return

        try:
            with open(HISTORICO_JSON_FILE, "r", encoding="utf-8") as f:
                historico = json.load(f)
        except:
            historico = []

        historico_valido = []
        alterou_json = False
        
        for h in historico:
            caminho_salvo = h.get("caminho", "")
            
            # GOLPE DA MIGRAÇÃO: Se a rota antiga não existir, redireciona para a nova pasta AppData
            if not os.path.exists(caminho_salvo):
                nome_arquivo = os.path.basename(caminho_salvo)
                caminho_tentativa = os.path.join(PASTA_HISTORICO_PPTX, nome_arquivo)
                if os.path.exists(caminho_tentativa):
                    h["caminho"] = caminho_tentativa
                    caminho_salvo = caminho_tentativa
                    alterou_json = True
            
            if os.path.exists(caminho_salvo):
                historico_valido.append(h)

        # Força limite de 5 na visualização (caso venha de versões antigas)
        historico_valido = historico_valido[-5:]

        # Salva as rotas consertadas de volta no arquivo
        if alterou_json:
            with open(HISTORICO_JSON_FILE, "w", encoding="utf-8") as f:
                json.dump(historico_valido, f, indent=4, ensure_ascii=False)

        if not historico_valido:
            QMessageBox.information(self, "Histórico Vazio", "Nenhuma apresentação física foi encontrada. Gere uma nova!")
            return

        # Monta a lista com no máximo 5 itens
        itens = [f"{item['data']} - {item.get('tema', 'Tema Desconhecido')}" for item in reversed(historico_valido)]
        
        item_selecionado, ok = QInputDialog.getItem(
            self, "Recuperar Apresentação Pronta", 
            "Escolha um PPTX já gerado anteriormente:", 
            itens, 0, False
        )

        if ok and item_selecionado:
            idx = itens.index(item_selecionado)
            caminho_escolhido = list(reversed(historico_valido))[idx]["caminho"]
            
            # MÁGICA: Carrega o PPTX nativo instantaneamente para o processador final
            self.pptx_path = caminho_escolhido
            self.lbl_pptx.setText(f"📁 Recuperado: {os.path.basename(caminho_escolhido)}")
            
            QMessageBox.information(self, "Sucesso Absoluto", "Apresentação recuperada instantaneamente!\n\nVocê já pode escolher suas fontes e cores e ir direto pro Passo 6 (Processar e Salvar).")

    def process_and_save(self):
        if not self.pptx_path:
            QMessageBox.warning(self, "Aviso", "Por favor, carregue uma apresentação primeiro.")
            return

        default_save_path = os.path.join(self.last_dir if self.last_dir else base_dir, "apresentacao_final.pptx")
        save_path, _ = QFileDialog.getSaveFileName(self, "Salvar Apresentação Final", default_save_path, "PowerPoint (*.pptx)")
        if not save_path: return
        
        self.last_dir = os.path.dirname(save_path)
        self.save_config()

        try:
            prs = Presentation(self.pptx_path)
            title_rgb = RGBColor(self.text_color_title.red(), self.text_color_title.green(), self.text_color_title.blue()) if self.text_color_title else None
            body_rgb = RGBColor(self.text_color_body.red(), self.text_color_body.green(), self.text_color_body.blue()) if self.text_color_body else None
            fill_rgb = RGBColor(self.fill_color.red(), self.fill_color.green(), self.fill_color.blue()) if self.fill_color else None

            for slide in prs.slides:
                # --- 1. CONTROLE DE LIGA/DESLIGA DA BARRA LATERAL (ITEM 1) ---
                shapes_para_remover = []
                for shape in slide.shapes:
                    # Identifica a barra tanto pela etiqueta nova quanto pelo formato antigo
                    is_barra = (shape.name == "BarraLateral_Vega") or (shape.left == 0 and shape.top == 0 and shape.width <= int(prs.slide_width * 0.03) and shape.height == prs.slide_height)
                    if not self.chk_barra_lateral.isChecked() and is_barra:
                        shapes_para_remover.append(shape)
                        
                for s_rem in shapes_para_remover:
                    slide.shapes._spTree.remove(s_rem._element)

                # --- 2. PROCESSAMENTO DAS OUTRAS FORMAS ---
                for shape in slide.shapes:
                    is_barra = (shape.name == "BarraLateral_Vega") or (shape.left == 0 and shape.top == 0 and shape.width <= int(prs.slide_width * 0.03) and shape.height == prs.slide_height)
                    if is_barra:
                        continue # Pula a barra lateral para não estragar a cor dela caso escolha cor de card!
                        
                    if self.chk_scale.isChecked() and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        scale = self.slider_scale.value() / 100.0 # <--- ESCALA BIDIRECIONAL (ITEM 2)
                        nw, nh = int(shape.width * scale), int(shape.height * scale)
                        shape.left += int((shape.width - nw) / 2)
                        shape.top += int((shape.height - nh) / 2)
                        shape.width, shape.height = nw, nh

                    if hasattr(shape, "fill") and fill_rgb:
                        try:
                            if shape.fill.type == MSO_FILL.SOLID: shape.fill.fore_color.rgb = fill_rgb
                        except: pass
                    
                    if shape.has_text_frame:
                        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Identifica quem é título e quem é corpo/card
                                is_title = (shape.name == "TitleBox_Vega")
                                if run.font.size:
                                    if run.font.size >= Pt(24): is_title = True
                                elif paragraph.font.size and paragraph.font.size >= Pt(24):
                                    is_title = True
                                    
                                # APLICA A COR SEPARADA (ITEM 3)
                                cor_rgb = title_rgb if is_title else body_rgb
                                if cor_rgb:
                                    run.font.color.rgb = cor_rgb
                                    
                                # APLICA A FONTE SEPARADA
                                fonte_escolhida = self.combo_font_title.currentText() if is_title else self.combo_font_body.currentText()
                                if fonte_escolhida != "Padrão do Tema":
                                    run.font.name = fonte_escolhida
                                
                                if self.chk_text_scale.isChecked():
                                    percentual = self.slider_text_scale.value()
                                    fator = 1.0 + (percentual / 100.0)
                                    tamanho = run.font.size if run.font.size is not None else paragraph.font.size
                                    if tamanho is None:
                                        tamanho = Pt(18)
                                    
                                    run.font.size = Pt(int(tamanho.pt * fator))

                if self.bg_path:
                    pic = slide.shapes.add_picture(self.bg_path, 0, 0, prs.slide_width, prs.slide_height)
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(2, pic._element)

            salvo_com_sucesso = False
            while not salvo_com_sucesso:
                try:
                    prs.save(save_path)
                    salvo_com_sucesso = True
                except PermissionError:
                    resp = QMessageBox.warning(
                        self, 
                        "Arquivo Aberto e Bloqueado", 
                        "O arquivo PowerPoint destino já está aberto e travado.\n\n"
                        "Deseja que o aplicativo feche o PowerPoint para continuar a edição?", 
                        QMessageBox.Ok | QMessageBox.Cancel
                    )
                    if resp == QMessageBox.Ok:
                        if os.name == 'nt':
                            os.system("taskkill /f /im POWERPNT.EXE")
                            import time; time.sleep(1.5)
                    else:
                        return 
                        
            # --- AUTO OPEN (Abre automaticamente o PPTX após edição) ---
            self.last_saved_pptx = save_path
            try:
                if os.name == 'nt': os.startfile(save_path)
                else: subprocess.Popen(['open' if sys.platform == 'darwin' else 'xdg-open', save_path])
            except: pass

        except Exception as e:
            QMessageBox.critical(self, "Erro", f"Erro ao processar arquivo:\n{str(e)}")
            
    def abrir_tutorial(self):
        import idiomas
        
        # Estrutura da Mágica Definitiva: (Widget Alvo, Título, Texto)
        passos = [
            (self.btn_config, idiomas.tr("tut_p1_tit"), idiomas.tr("tut_p1_txt")),
            (self.btn_novidades, idiomas.tr("tut_p2_tit"), idiomas.tr("tut_p2_txt")),
            (self.btn_load, idiomas.tr("tut_p3_tit"), idiomas.tr("tut_p3_txt")),
            (self.txt_tema, idiomas.tr("tut_p4_tit"), idiomas.tr("tut_p4_txt")),
            (None, idiomas.tr("tut_p5_tit"), idiomas.tr("tut_p5_txt")), # Janela de Imagens
            (self.btn_historico, idiomas.tr("tut_p6_tit"), idiomas.tr("tut_p6_txt")),
            (self.combo_presets, idiomas.tr("tut_p7_tit"), idiomas.tr("tut_p7_txt")),
            (self.btn_color_title, idiomas.tr("tut_p8_tit"), idiomas.tr("tut_p8_txt")),
            (self.combo_font_title, idiomas.tr("tut_p9_tit"), idiomas.tr("tut_p9_txt")),
            (self.btn_color_body, idiomas.tr("tut_p10_tit"), idiomas.tr("tut_p10_txt")),
            (self.combo_font_body, idiomas.tr("tut_p11_tit"), idiomas.tr("tut_p11_txt")),
            (self.btn_fill_color, idiomas.tr("tut_p12_tit"), idiomas.tr("tut_p12_txt")),
            (self.btn_bg, idiomas.tr("tut_p13_tit"), idiomas.tr("tut_p13_txt")),
            (self.chk_scale, idiomas.tr("tut_p14_tit"), idiomas.tr("tut_p14_txt")),
            (self.chk_text_scale, idiomas.tr("tut_p15_tit"), idiomas.tr("tut_p15_txt")),
            (self.frame_preview, idiomas.tr("tut_p16_tit"), idiomas.tr("tut_p16_txt")),
            (self.chk_notes, idiomas.tr("tut_p17_tit"), idiomas.tr("tut_p17_txt")),
            (self.btn_save, idiomas.tr("tut_p18_tit"), idiomas.tr("tut_p18_txt")),
            (self.btn_pdf, idiomas.tr("tut_p19_tit"), idiomas.tr("tut_p19_txt"))
        ]
        
        self.overlay_tutorial.iniciar(passos)

    def resizeEvent(self, event):
        super().resizeEvent(event)
        # Garante que o recorte cubra a tela toda se o cara redimensionar a janela durante o tutorial
        if hasattr(self, 'overlay_tutorial') and self.overlay_tutorial.isVisible():
            self.overlay_tutorial.setGeometry(0, 0, self.width(), self.height())
            self.overlay_tutorial.atualizar_passo()
            
    def checar_notas_atualizacao(self):
        # Olha no config.json se ele já viu as novidades
        ultima_vista = self.config.get("ultima_versao_vista", "")
        
        try:
            # Puxa a versão correta do servidor em vez do manifesto local desatualizado
            url = f"https://vegap-masterapp.hf.space/master/atualizacao/{PRODUTO_ID_MASTER}"
            resposta = requests.get(url, timeout=3)
            versao_real = resposta.json().get("versao_atual", VERSAO_ATUAL) if resposta.status_code == 200 else VERSAO_ATUAL
        except:
            versao_real = VERSAO_ATUAL
            
        if ultima_vista != versao_real:
            import idiomas
            msg = QMessageBox(self)
            msg.setWindowTitle(idiomas.tr("msg_update_titulo").format(versao_real))
            msg.setText(idiomas.tr("msg_update_texto").format(versao_real))
            msg.setIcon(QMessageBox.Information)
            msg.setStandardButtons(QMessageBox.Ok)
            msg.exec()
            
            self.config["ultima_versao_vista"] = versao_real
            self.save_config()
            
    def abrir_novidades(self):
        import idiomas
        # BUSCA AO VIVO NO BACKEND DO MASTER
        self.btn_novidades.setText(idiomas.tr("btn_buscando"))
        self.btn_novidades.setEnabled(False)
        QApplication.processEvents()
        
        try:
            url = f"https://vegap-masterapp.hf.space/master/atualizacao/{PRODUTO_ID_MASTER}"
            resposta = requests.get(url, timeout=5)
            versao_nuvem = VERSAO_ATUAL
            if resposta.status_code == 200:
                dados = resposta.json()
                
                import idiomas
                if idiomas._idioma_atual == "en":
                    notas = dados.get("notas_atualizacao_en", "")
                else:
                    notas = dados.get("notas_atualizacao", "")
                    
                notas = notas.strip() if notas else ""
                
                # PEGA A VERSÃO REAL DA API PARA O TÍTULO BATER COM O TEXTO!
                versao_nuvem = dados.get("versao_atual", VERSAO_ATUAL)
                if not notas:
                    notas = idiomas.tr("msg_update_vazio")
            else:
                import idiomas
                notas = idiomas.tr("msg_update_erro_api")
        except Exception as e:
            import idiomas
            notas = idiomas.tr("msg_update_erro_con").format(e)
        finally:
            import idiomas
            self.btn_novidades.setText(idiomas.tr("btn_novidades"))
            self.btn_novidades.setEnabled(True)

        # ========================================================
        # FORMATADOR INTELIGENTE (Puxa do Banco e Quebra Linhas)
        # ========================================================
        notas_formatadas = notas.replace("\\n", "\n") # Desescapa possíveis quebras da API
        notas_formatadas = notas_formatadas.replace(" • ", "\n• ")
        
        # Dá um respiro antes de cada emoji de categoria
        for emj in ["🎨", "🛡️", "⚡", "🚀", "♿", "💡", "✨", "💳", "🔧", "🔥"]:
            notas_formatadas = notas_formatadas.replace(f" {emj}", f"\n\n{emj}")
            
        # Limpa as quebras de linha excessivas geradas pelo replace
        notas_formatadas = "\n".join([linha.strip() for linha in notas_formatadas.split("\n") if linha.strip()])

        zoom_atual = self.config.get("zoom", 100)
        # USA A VERSÃO DA NUVEM PARA A JANELA
        dialog = NovidadesDialog(versao_nuvem, notas_formatadas, zoom_atual, self)
        dialog.exec()

    def export_to_pdf(self):
        # Tenta pegar o último que o usuário editou/salvou. Se não tiver, pega o que ele carregou no passo 1.
        alvo = self.last_saved_pptx if self.last_saved_pptx else self.pptx_path
        
        if not alvo or not os.path.exists(alvo):
            QMessageBox.warning(self, "Aviso", "Por favor, carregue ou salve uma apresentação primeiro.")
            return
            
        pdf_path, _ = QFileDialog.getSaveFileName(self, "Salvar PDF", alvo.replace(".pptx", ".pdf"), "PDF (*.pdf)")
        if not pdf_path: return
        
        progress = QProgressDialog("Convertendo para PDF (Aguarde...)", "Cancelar", 0, 0, self)
        progress.setWindowTitle("Gerando PDF")
        progress.setWindowModality(Qt.WindowModal)
        progress.show()
        QApplication.processEvents()
        
        try:
            # 32 é o código do formato PDF no PowerPoint
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(os.path.abspath(alvo))
            deck.SaveAs(os.path.abspath(pdf_path), 32)
            deck.Close()
            powerpoint.Quit()
            
            progress.close()
            QMessageBox.information(self, "Sucesso", "Apresentação convertida para PDF com sucesso!")
            
            if os.name == 'nt': os.startfile(pdf_path)
        except Exception as e:
            progress.close()
            QMessageBox.critical(self, "Erro no PDF", f"Falha ao converter.\nVerifique se o PowerPoint não está com janelas de diálogo travando o fundo.\n\nDetalhes:\n{str(e)}")
            
if __name__ == "__main__":
    if splash: splash.atualizar(85, "Verificando radares e atualizações do sistema...")
    if checar_e_atualizar():
        os._exit(0)
        
    if splash: splash.atualizar(95, "Preparando interface gráfica e blindagens...")
    window = AppPaiVega()
    
    if splash:
        splash.atualizar(100, "Tudo pronto! Bem-vindo(a) ao palco.")
        time.sleep(0.5)
        splash.finish(window)
        
    window.show()
    sys.exit(app.exec())